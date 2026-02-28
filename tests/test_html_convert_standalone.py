# -*- coding: utf-8 -*-
"""
HTML to PPTX 独立测试脚本
"""

import os
import re
import base64
import urllib.request
from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path

from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class CSSParser:
    """CSS样式解析器"""

    @staticmethod
    def parse_inline_style(style_str: str) -> Dict[str, str]:
        if not style_str:
            return {}
        styles = {}
        declarations = style_str.split(';')
        for decl in declarations:
            decl = decl.strip()
            if ':' in decl:
                prop, value = decl.split(':', 1)
                styles[prop.strip().lower()] = value.strip()
        return styles

    @staticmethod
    def parse_length(value: str, reference: float = 1920) -> float:
        if not value:
            return 0
        value = value.strip().lower()
        if value.isdigit():
            return float(value)
        match = re.match(r'^(-?\d*\.?\d+)(px|pt|in|cm|mm|em|rem|%|vw|vh)?$', value)
        if match:
            num = float(match.group(1))
            unit = match.group(2) or 'px'
            if unit == 'px':
                return num
            elif unit == 'pt':
                return num * 1.333
            elif unit == 'in':
                return num * 96
            elif unit == 'cm':
                return num * 37.795
            elif unit == 'mm':
                return num * 3.7795
            elif unit in ('em', 'rem'):
                return num * 16
            elif unit == '%':
                return num / 100 * reference
            elif unit == 'vw':
                return num / 100 * 1920
            elif unit == 'vh':
                return num / 100 * 1080
        return 0

    @staticmethod
    def parse_color(value: str) -> Optional[Tuple[int, int, int]]:
        if not value:
            return None
        value = value.strip().lower()
        if value == 'transparent':
            return None
        if value.startswith('#'):
            hex_color = value[1:]
            if len(hex_color) == 3:
                hex_color = ''.join([c*2 for c in hex_color])
            if len(hex_color) == 6:
                try:
                    r = int(hex_color[0:2], 16)
                    g = int(hex_color[2:4], 16)
                    b = int(hex_color[4:6], 16)
                    return (r, g, b)
                except ValueError:
                    pass
        rgb_match = re.match(r'rgba?\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)', value)
        if rgb_match:
            return (int(rgb_match.group(1)), int(rgb_match.group(2)), int(rgb_match.group(3)))
        color_names = {
            'white': (255, 255, 255), 'black': (0, 0, 0), 'red': (255, 0, 0),
            'green': (0, 128, 0), 'blue': (0, 0, 255), 'yellow': (255, 255, 0),
        }
        if value in color_names:
            return color_names[value]
        return None

    @staticmethod
    def parse_gradient(value: str) -> Optional[Dict[str, Any]]:
        if not value:
            return None
        linear_match = re.match(r'linear-gradient\s*\(\s*(.+)\s*\)', value, re.IGNORECASE)
        if linear_match:
            content = linear_match.group(1)
            parts = []
            depth = 0
            current = ""
            for char in content:
                if char == '(':
                    depth += 1
                elif char == ')':
                    depth -= 1
                elif char == ',' and depth == 0:
                    parts.append(current.strip())
                    current = ""
                    continue
                current += char
            if current.strip():
                parts.append(current.strip())

            angle = 180
            stops = []
            for i, part in enumerate(parts):
                part = part.strip()
                if i == 0:
                    angle_match = re.match(r'(\d+)deg', part)
                    if angle_match:
                        angle = int(angle_match.group(1))
                        continue
                    if 'to right' in part:
                        angle = 90
                        continue
                    elif 'to left' in part:
                        angle = 270
                        continue
                    elif 'to bottom' in part:
                        angle = 180
                        continue
                    elif 'to top' in part:
                        angle = 0
                        continue
                color = CSSParser.parse_color(part.split()[0])
                if color:
                    pos_match = re.search(r'(\d+)%', part)
                    pos = int(pos_match.group(1)) if pos_match else None
                    stops.append({'color': color, 'position': pos})
            if stops:
                return {'type': 'linear', 'angle': angle, 'stops': stops}
        return None


class HTMLToPPTXConverter:
    """HTML到PPTX转换器"""

    def __init__(self, html_content: str, base_path: str = None):
        self.soup = BeautifulSoup(html_content, 'html.parser')
        self.base_path = base_path or os.getcwd()
        self.css_parser = CSSParser()
        self.slide_width = 1920
        self.slide_height = 1080
        self.ppt_width_inches = 13.333
        self.ppt_height_inches = 7.5
        self.scale_x = self.ppt_width_inches / self.slide_width
        self.scale_y = self.ppt_height_inches / self.slide_height
        self.temp_images = []

    def px_to_inches(self, px: float, axis: str = 'x') -> float:
        scale = self.scale_x if axis == 'x' else self.scale_y
        return px * scale

    def detect_slide_size(self):
        body = self.soup.find('body')
        if body:
            for child in body.children:
                if isinstance(child, Tag):
                    styles = self.css_parser.parse_inline_style(child.get('style', ''))
                    width = styles.get('width', '')
                    height = styles.get('height', '')
                    if width and height:
                        w = self.css_parser.parse_length(width)
                        h = self.css_parser.parse_length(height)
                        if w > 0 and h > 0:
                            self.slide_width = w
                            self.slide_height = h
                            self.scale_x = self.ppt_width_inches / self.slide_width
                            self.scale_y = self.ppt_height_inches / self.slide_height
                            print(f"  检测到幻灯片尺寸: {w}x{h}")
                            return

    def get_element_bounds(self, element: Tag, parent_bounds: Dict = None) -> Dict[str, float]:
        styles = self.css_parser.parse_inline_style(element.get('style', ''))
        parent_bounds = parent_bounds or {
            'left': 0, 'top': 0, 'width': self.slide_width, 'height': self.slide_height
        }
        position = styles.get('position', 'static')
        width = self.css_parser.parse_length(styles.get('width', ''), parent_bounds['width'])
        height = self.css_parser.parse_length(styles.get('height', ''), parent_bounds['height'])

        if position == 'absolute':
            left = self.css_parser.parse_length(styles.get('left', '0'), parent_bounds['width'])
            top = self.css_parser.parse_length(styles.get('top', '0'), parent_bounds['height'])
            if 'right' in styles and 'left' not in styles:
                right = self.css_parser.parse_length(styles['right'], parent_bounds['width'])
                left = parent_bounds['width'] - right - width
            if 'bottom' in styles and 'top' not in styles:
                bottom = self.css_parser.parse_length(styles['bottom'], parent_bounds['height'])
                top = parent_bounds['height'] - bottom - height
        else:
            left = parent_bounds.get('left', 0)
            top = parent_bounds.get('top', 0)

        return {
            'left': left, 'top': top,
            'width': width if width > 0 else parent_bounds['width'],
            'height': height if height > 0 else 50
        }

    def apply_shape_fill(self, shape, styles: Dict[str, str]):
        background = styles.get('background', '')
        bg_color = styles.get('background-color', '')

        gradient = self.css_parser.parse_gradient(background) or self.css_parser.parse_gradient(bg_color)
        if gradient and gradient['stops']:
            try:
                fill = shape.fill
                fill.gradient()
                fill.gradient_angle = gradient.get('angle', 180)
                fill.gradient_stops[0].color.rgb = RGBColor(*gradient['stops'][0]['color'])
                if len(gradient['stops']) > 1:
                    fill.gradient_stops[1].color.rgb = RGBColor(*gradient['stops'][-1]['color'])
                return {'type': 'gradient'}
            except Exception as e:
                print(f"  渐变填充失败: {e}")

        color = self.css_parser.parse_color(bg_color) or self.css_parser.parse_color(background)
        if color:
            try:
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*color)
                return {'type': 'solid'}
            except Exception as e:
                print(f"  纯色填充失败: {e}")
        return None

    def apply_text_format(self, paragraph, run, styles: Dict[str, str]):
        font = run.font
        font_size = styles.get('font-size', '')
        if font_size:
            size_px = self.css_parser.parse_length(font_size)
            if size_px > 0:
                font.size = Pt(size_px * 0.75)
        color = self.css_parser.parse_color(styles.get('color', ''))
        if color:
            font.color.rgb = RGBColor(*color)
        font_weight = styles.get('font-weight', '')
        if font_weight in ('bold', '700', '800', '900'):
            font.bold = True
        if styles.get('font-style') == 'italic':
            font.italic = True
        font_family = styles.get('font-family', '')
        if font_family:
            first_font = font_family.split(',')[0].strip().strip("'\"")
            font.name = first_font
        text_align = styles.get('text-align', '')
        if text_align == 'center':
            paragraph.alignment = PP_ALIGN.CENTER
        elif text_align == 'right':
            paragraph.alignment = PP_ALIGN.RIGHT

    def process_element(self, element: Tag, slide, parent_bounds: Dict = None, z_index: int = 0) -> int:
        if not isinstance(element, Tag):
            return z_index

        styles = self.css_parser.parse_inline_style(element.get('style', ''))
        bounds = self.get_element_bounds(element, parent_bounds)
        elem_z = int(styles.get('z-index', z_index))

        has_background = any(k in styles for k in ['background', 'background-color', 'background-image'])
        has_border = 'border' in styles

        direct_text = ''
        if element.name == 'p':
            direct_text = element.get_text(strip=True)

        left = Inches(self.px_to_inches(bounds['left'], 'x'))
        top = Inches(self.px_to_inches(bounds['top'], 'y'))
        width = Inches(self.px_to_inches(bounds['width'], 'x'))
        height = Inches(self.px_to_inches(bounds['height'], 'y'))

        shape = None

        if has_background or has_border:
            border_radius = styles.get('border-radius', '')
            radius_px = self.css_parser.parse_length(border_radius) if border_radius else 0

            if radius_px > 20:
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
                try:
                    shape.adjustments[0] = min(radius_px / min(bounds['width'], bounds['height']), 0.5)
                except (IndexError, TypeError):
                    pass
            else:
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

            self.apply_shape_fill(shape, styles)

            if shape and hasattr(shape, 'line'):
                shape.line.fill.background()

        if direct_text and element.name == 'p':
            if shape:
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = direct_text
                self.apply_text_format(p, run, styles)
                tf.anchor = MSO_ANCHOR.MIDDLE
            else:
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = direct_text
                self.apply_text_format(p, run, styles)

        next_z = elem_z + 1
        for child in element.children:
            if isinstance(child, Tag):
                child_parent_bounds = {
                    'left': bounds['left'], 'top': bounds['top'],
                    'width': bounds['width'], 'height': bounds['height']
                }
                next_z = self.process_element(child, slide, child_parent_bounds, next_z)

        return next_z

    def convert(self, output_path: str) -> Dict[str, Any]:
        try:
            self.detect_slide_size()

            prs = Presentation()
            prs.slide_width = Inches(self.ppt_width_inches)
            prs.slide_height = Inches(self.ppt_height_inches)

            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            body = self.soup.find('body')
            if body:
                main_container = None
                for child in body.children:
                    if isinstance(child, Tag):
                        styles = self.css_parser.parse_inline_style(child.get('style', ''))
                        if styles.get('width') and styles.get('height'):
                            main_container = child
                            break

                if main_container:
                    container_styles = self.css_parser.parse_inline_style(main_container.get('style', ''))

                    bg_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        Inches(self.ppt_width_inches), Inches(self.ppt_height_inches)
                    )
                    bg_shape.line.fill.background()
                    self.apply_shape_fill(bg_shape, container_styles)

                    container_bounds = {
                        'left': 0, 'top': 0,
                        'width': self.slide_width, 'height': self.slide_height
                    }

                    element_count = 0
                    for child in main_container.children:
                        if isinstance(child, Tag):
                            self.process_element(child, slide, container_bounds)
                            element_count += 1

                    print(f"  处理了 {element_count} 个顶层元素")

            prs.save(output_path)

            return {
                'success': True,
                'output_path': output_path,
                'slide_size': f'{self.slide_width}x{self.slide_height}',
                'message': 'HTML successfully converted to PPTX'
            }

        except Exception as e:
            import traceback
            traceback.print_exc()
            return {'success': False, 'error': str(e)}


# 测试
if __name__ == "__main__":
    print("=== HTML to PPTX 转换测试 ===\n")

    html_path = 'test_output/slide_page.html'
    output_path = 'test_output/html_converted.pptx'

    print(f"读取 HTML: {html_path}")
    with open(html_path, 'r', encoding='utf-8') as f:
        html_content = f.read()

    print("开始转换...")
    converter = HTMLToPPTXConverter(html_content, os.path.dirname(html_path))
    result = converter.convert(output_path)
    assert isinstance(result, dict), "Expected dict result"

    print(f"\n转换结果: {result}")
