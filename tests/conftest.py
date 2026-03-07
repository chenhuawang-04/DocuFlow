"""DocuFlow test configuration and shared fixtures."""

import os
import shutil
import tempfile
from pathlib import Path

import pytest


# ---------------------------------------------------------------------------
# Skip markers for external tool dependencies
# ---------------------------------------------------------------------------

requires_tesseract = pytest.mark.skipif(
    shutil.which("tesseract") is None,
    reason="Tesseract OCR not installed",
)

requires_pandoc = pytest.mark.skipif(
    shutil.which("pandoc") is None,
    reason="pandoc not installed",
)






# ---------------------------------------------------------------------------
# Shared fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def tmp_dir(tmp_path):
    """Provide a temporary directory that is cleaned up after the test."""
    return tmp_path


@pytest.fixture
def sample_docx(tmp_path):
    """Create a minimal .docx file for testing."""
    try:
        from docx import Document
    except ImportError:
        pytest.skip("python-docx not installed")
    doc = Document()
    doc.add_paragraph("Test paragraph")
    path = tmp_path / "sample.docx"
    doc.save(str(path))
    return str(path)


@pytest.fixture
def sample_xlsx(tmp_path):
    """Create a minimal .xlsx file for testing."""
    try:
        from openpyxl import Workbook
    except ImportError:
        pytest.skip("openpyxl not installed")
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Hello"
    ws["B1"] = 42
    path = tmp_path / "sample.xlsx"
    wb.save(str(path))
    wb.close()
    return str(path)


@pytest.fixture
def sample_pptx(tmp_path):
    """Create a minimal .pptx file for testing."""
    try:
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not installed")
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    return str(path)
