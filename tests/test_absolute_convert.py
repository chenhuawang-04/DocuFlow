# -*- coding: utf-8 -*-
"""HTML绝对定位转PPTX - 修复版"""
import os
import re
from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


def make_oxml_element(tag):
    """创建 OXML 元素"""
    return etree.Element(qn(tag))


class CSSParser:
    @staticmethod
    def parse_inline_style(style_str):
        if not style_str:
            return {}
        styles = {}
        for decl in style_str.split(';'):
            decl = decl.strip()
            if ':' in decl:
                prop, value = decl.split(':', 1)
                styles[prop.strip().lower()] = value.strip()
        return styles

    @staticmethod
    def parse_length(value, reference=1920):
        if not value:
            return 0
        value = value.strip().lower()
        match = re.match(r'^(-?\d*\.?\d+)(px|pt|in|cm|%)?$', value)
        if match:
            num = float(match.group(1))
            unit = match.group(2) or 'px'
            if unit == 'px':
                return num
            elif unit == 'in':
                return num * 96
            elif unit == '%':
                return num / 100 * reference
        return 0

    @staticmethod
    def parse_color(value):
        """解析颜色，支持 #hex, rgb(), rgba()，返回 (r, g, b, alpha)"""
        if not value:
            return None
        value = value.strip().lower()

        # #hex
        if value.startswith('#'):
            h = value[1:]
            if len(h) == 3:
                h = ''.join([c * 2 for c in h])
            if len(h) == 6:
                return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16), 1.0)

        # rgba(r, g, b, a)
        m = re.match(r'rgba\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*,\s*([\d.]+)\s*\)', value)
        if m:
            return (int(m.group(1)), int(m.group(2)), int(m.group(3)), float(m.group(4)))

        # rgb(r, g, b)
        m = re.match(r'rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', value)
        if m:
            return (int(m.group(1)), int(m.group(2)), int(m.group(3)), 1.0)

        # 颜色名
        colors = {
            'white': (255, 255, 255, 1.0),
            'black': (0, 0, 0, 1.0),
            'red': (255, 0, 0, 1.0),
            'green': (0, 128, 0, 1.0),
            'blue': (0, 0, 255, 1.0),
        }
        if value in colors:
            return colors[value]

        return None

    @staticmethod
    def parse_gradient(value):
        if not value:
            return None
        m = re.match(r'linear-gradient\s*\(\s*(.+)\s*\)', value, re.IGNORECASE)
        if m:
            content = m.group(1)
            parts = [p.strip() for p in re.split(r',(?![^()]*\))', content)]
            angle = 180
            stops = []
            for i, part in enumerate(parts):
                if i == 0:
                    am = re.match(r'(\d+)deg', part)
                    if am:
                        angle = int(am.group(1))
                        continue
                color = CSSParser.parse_color(part.split()[0])
                if color:
                    stops.append(color[:3])  # 只取RGB，忽略alpha
            if stops:
                return {'angle': angle, 'stops': stops}
        return None


# 转换参数
SLIDE_W, SLIDE_H = 1920, 1080
PPT_W, PPT_H = 13.333, 7.5
SCALE_X, SCALE_Y = PPT_W / SLIDE_W, PPT_H / SLIDE_H

# 字体缩放：px -> pt，同时考虑画布缩放
# 1920px 对应 13.333in，1in = 72pt
# font_pt = font_px * (PPT_W * 72 / SLIDE_W) = font_px * 0.5
FONT_SCALE = PPT_W * 72 / SLIDE_W  # ≈ 0.5


def px_to_in(px, axis='x'):
    return px * (SCALE_X if axis == 'x' else SCALE_Y)


def set_shape_rounded_corners(shape, radius_inches):
    """设置形状的圆角"""
    # 通过XML设置圆角
    spPr = shape._element.spPr
    prstGeom = spPr.prstGeom
    if prstGeom is not None:
        # 设置为圆角矩形
        prstGeom.set('prst', 'roundRect')
        # 添加调整值
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            from lxml import etree
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        # 圆角大小 (0-50000, 表示百分比*1000)
        # 计算圆角比例
        min_dim = min(shape.width.inches, shape.height.inches)
        if min_dim > 0:
            ratio = min(radius_inches / min_dim * 50000, 50000)
            from lxml import etree
            # 清除旧的
            for child in list(avLst):
                avLst.remove(child)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {int(ratio)}')


def convert_html_to_pptx(html_path, output_path):
    # 读取HTML
    with open(html_path, 'r', encoding='utf-8') as f:
        html = f.read()

    soup = BeautifulSoup(html, 'html.parser')
    css = CSSParser()

    # 创建PPT
    prs = Presentation()
    prs.slide_width = Inches(PPT_W)
    prs.slide_height = Inches(PPT_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 找主容器
    body = soup.find('body')
    main = None
    for child in body.children:
        if isinstance(child, Tag):
            s = css.parse_inline_style(child.get('style', ''))
            if s.get('width') and s.get('height'):
                main = child
                break

    if main:
        # 背景
        ms = css.parse_inline_style(main.get('style', ''))
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(PPT_W), Inches(PPT_H))
        bg.line.fill.background()
        grad = css.parse_gradient(ms.get('background', ''))
        if grad and grad['stops']:
            bg.fill.gradient()
            bg.fill.gradient_angle = grad['angle']
            bg.fill.gradient_stops[0].color.rgb = RGBColor(*grad['stops'][0])
            if len(grad['stops']) > 1:
                bg.fill.gradient_stops[1].color.rgb = RGBColor(*grad['stops'][-1])

        # 处理子元素
        count = 0
        for elem in main.descendants:
            if not isinstance(elem, Tag):
                continue
            s = css.parse_inline_style(elem.get('style', ''))
            if s.get('position') != 'absolute':
                continue

            # 解析位置和大小
            left = css.parse_length(s.get('left', '0'))
            top = css.parse_length(s.get('top', '0'))
            w = css.parse_length(s.get('width', '100'))
            h = css.parse_length(s.get('height', '50'))

            # 处理 bottom/right
            if 'bottom' in s and 'top' not in s:
                bottom = css.parse_length(s['bottom'])
                top = SLIDE_H - bottom - h
            if 'right' in s and 'left' not in s:
                right = css.parse_length(s['right'])
                left = SLIDE_W - right - w

            # 转换为英寸
            left_in = px_to_in(left, 'x')
            top_in = px_to_in(top, 'y')
            w_in = px_to_in(w, 'x')
            h_in = px_to_in(h, 'y')

            L, T = Inches(left_in), Inches(top_in)
            W, H = Inches(w_in), Inches(h_in)

            # 创建形状
            if elem.name == 'div':
                radius_px = css.parse_length(s.get('border-radius', '0'))
                radius_in = px_to_in(radius_px, 'x')

                # 先创建矩形
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, L, T, W, H)

                # 如果有圆角，设置圆角
                if radius_px > 0:
                    set_shape_rounded_corners(shape, radius_in)

                shape.line.fill.background()

                # 填充
                grad = css.parse_gradient(s.get('background', ''))
                if grad and grad['stops']:
                    shape.fill.gradient()
                    shape.fill.gradient_angle = grad['angle']
                    shape.fill.gradient_stops[0].color.rgb = RGBColor(*grad['stops'][0])
                    if len(grad['stops']) > 1:
                        shape.fill.gradient_stops[1].color.rgb = RGBColor(*grad['stops'][-1])
                else:
                    bg_str = s.get('background', '') or s.get('background-color', '')
                    color = css.parse_color(bg_str)
                    if color:
                        r, g, b, alpha = color
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(r, g, b)
                        # 设置透明度
                        if alpha < 1.0:
                            spPr = shape._element.spPr
                            solidFill = spPr.find(qn('a:solidFill'))
                            if solidFill is not None:
                                srgbClr = solidFill.find(qn('a:srgbClr'))
                                if srgbClr is not None:
                                    alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
                                    alpha_elem.set('val', str(int(alpha * 100000)))
                count += 1

            elif elem.name == 'p':
                text = elem.get_text(strip=True)
                if text:
                    tb = slide.shapes.add_textbox(L, T, W, H)
                    tf = tb.text_frame
                    tf.word_wrap = False
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = text

                    # 字体
                    font = run.font

                    # 字体大小 - 按画布比例缩放
                    fs = css.parse_length(s.get('font-size', '16'))
                    if fs > 0:
                        font.size = Pt(fs * FONT_SCALE)

                    # 颜色 - 支持rgba透明度
                    color_str = s.get('color', '')
                    color = css.parse_color(color_str)
                    if color:
                        r, g, b, alpha = color
                        font.color.rgb = RGBColor(r, g, b)

                        # 设置透明度 (alpha < 1.0)
                        if alpha < 1.0:
                            # 透明度 = 1 - opacity，PPT用百分比*1000
                            transparency = int((1 - alpha) * 100000)
                            # 通过XML设置
                            rPr = run._r.get_or_add_rPr()
                            solidFill = rPr.find(qn('a:solidFill'))
                            if solidFill is not None:
                                srgbClr = solidFill.find(qn('a:srgbClr'))
                                if srgbClr is not None:
                                    # 移除旧的alpha
                                    for old_alpha in srgbClr.findall(qn('a:alpha')):
                                        srgbClr.remove(old_alpha)
                                    # 添加新alpha
                                    alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
                                    alpha_elem.set('val', str(int(alpha * 100000)))

                        print(f'  Text "{text[:20]}" color: {color_str} -> RGB({r},{g},{b}) alpha={alpha}')

                    fw = s.get('font-weight', '')
                    if fw in ('bold', '700', '800', '900'):
                        font.bold = True
                    ff = s.get('font-family', '')
                    if ff:
                        font.name = ff.split(',')[0].strip().strip("'\"")

                    # 对齐
                    ta = s.get('text-align', '')
                    if ta == 'center':
                        p.alignment = PP_ALIGN.CENTER
                    elif ta == 'right':
                        p.alignment = PP_ALIGN.RIGHT

                    count += 1

        print(f'Processed {count} elements')

    prs.save(output_path)
    print(f'Saved: {output_path}')
    return output_path


if __name__ == '__main__':
    convert_html_to_pptx('test_output/slide_absolute.html', 'test_output/slide_absolute.pptx')
