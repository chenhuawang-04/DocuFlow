"""
DocuFlow PPT - PowerPoint文档操作模块

提供PPT文档的创建、读取、编辑等功能
依赖: python-pptx
"""
import os
import io
from pathlib import Path
from typing import Optional, List, Dict, Any, Union

from ..core.registry import register_tool
from ..utils.deps import check_import


def _parse_length(value: str) -> int:
    """解析长度值，支持 'Xin', 'Xcm', 'Xpt', 'Xemu' 格式"""
    if not check_import("pptx"):
        return 0

    from pptx.util import Inches, Cm, Pt, Emu

    if isinstance(value, (int, float)):
        return Inches(value)

    value = str(value).strip().lower()

    if value.endswith('in'):
        return Inches(float(value[:-2]))
    elif value.endswith('cm'):
        return Cm(float(value[:-2]))
    elif value.endswith('pt'):
        return Pt(float(value[:-2]))
    elif value.endswith('emu'):
        return int(float(value[:-3]))
    else:
        # 默认当作英寸
        try:
            return Inches(float(value))
        except ValueError:
            return Inches(1)


class PPTOperations:
    """PowerPoint文档操作"""

    # ========== 文档操作 ==========

    @register_tool("ppt_create",
                   required_params=['path'],
                   optional_params=['title', 'width', 'height'])
    @staticmethod
    def create(path: str,
               title: Optional[str] = None,
               width: Optional[str] = None,
               height: Optional[str] = None) -> Dict[str, Any]:
        """
        创建新的PowerPoint文档

        Args:
            path: 文件保存路径，必须以.pptx结尾
            title: 可选的文档标题（元数据）
            width: 幻灯片宽度，如 '10in', '25.4cm'
            height: 幻灯片高度，如 '7.5in', '19.05cm'

        Returns:
            {success, path, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from pptx.util import Inches

            if not path.endswith('.pptx'):
                return {"success": False, "error": "文件路径必须以.pptx结尾"}

            # 创建演示文稿
            prs = Presentation()

            # 设置幻灯片尺寸
            if width:
                prs.slide_width = _parse_length(width)
            if height:
                prs.slide_height = _parse_length(height)

            # 设置文档属性
            if title:
                prs.core_properties.title = title

            # 创建目录
            dir_path = os.path.dirname(path)
            if dir_path:
                os.makedirs(dir_path, exist_ok=True)
            # 保存文件
            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height,
                "message": f"已创建PPT文档: {path}"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("ppt_read",
                   required_params=['path'],
                   optional_params=['include_notes'])
    @staticmethod
    def read(path: str,
             include_notes: bool = False) -> Dict[str, Any]:
        """
        读取PowerPoint文档内容

        Args:
            path: PPT文件路径
            include_notes: 是否包含演讲者备注

        Returns:
            {success, slides, total_slides, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)
            slides_data = []

            for idx, slide in enumerate(prs.slides):
                slide_info = {
                    "index": idx + 1,
                    "layout": slide.slide_layout.name if slide.slide_layout else "Unknown",
                    "shapes": [],
                    "text_content": []
                }

                # 提取形状信息
                for shape in slide.shapes:
                    shape_info = {
                        "name": shape.name,
                        "type": str(shape.shape_type),
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }

                    # 提取文本
                    if shape.has_text_frame:
                        text = ""
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text += run.text
                            text += "\n"
                        text = text.strip()
                        if text:
                            shape_info["text"] = text
                            slide_info["text_content"].append(text)

                    # 提取表格
                    if shape.has_table:
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text)
                            table_data.append(row_data)
                        shape_info["table"] = table_data

                    slide_info["shapes"].append(shape_info)

                # 提取备注
                if include_notes and slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    notes_text = notes_slide.notes_text_frame.text
                    slide_info["notes"] = notes_text

                slides_data.append(slide_info)

            return {
                "success": True,
                "path": path,
                "total_slides": len(prs.slides),
                "slides": slides_data,
                "message": f"已读取 {len(prs.slides)} 张幻灯片"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("ppt_info",
                   required_params=['path'],
                   optional_params=[])
    @staticmethod
    def info(path: str) -> Dict[str, Any]:
        """
        获取PowerPoint文档基本信息

        Args:
            path: PPT文件路径

        Returns:
            {success, slides, properties, dimensions, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from pptx.util import Inches

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            # 获取文档属性
            props = prs.core_properties
            properties = {
                "title": props.title,
                "author": props.author,
                "subject": props.subject,
                "keywords": props.keywords,
                "comments": props.comments,
                "created": str(props.created) if props.created else None,
                "modified": str(props.modified) if props.modified else None,
                "last_modified_by": props.last_modified_by
            }

            # 获取尺寸信息
            dimensions = {
                "width_emu": prs.slide_width,
                "height_emu": prs.slide_height,
                "width_inches": round(prs.slide_width / 914400, 2),
                "height_inches": round(prs.slide_height / 914400, 2)
            }

            # 获取布局统计
            layout_counts = {}
            for slide in prs.slides:
                layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
                layout_counts[layout_name] = layout_counts.get(layout_name, 0) + 1

            # 文件大小
            file_size = os.path.getsize(path)

            return {
                "success": True,
                "path": path,
                "total_slides": len(prs.slides),
                "properties": properties,
                "dimensions": dimensions,
                "layout_counts": layout_counts,
                "file_size_bytes": file_size,
                "file_size_mb": round(file_size / 1024 / 1024, 2),
                "message": f"PPT文档包含 {len(prs.slides)} 张幻灯片"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("ppt_set_properties",
                   required_params=['path'],
                   optional_params=['title', 'author', 'subject', 'keywords', 'comments'])
    @staticmethod
    def set_properties(path: str,
                       title: Optional[str] = None,
                       author: Optional[str] = None,
                       subject: Optional[str] = None,
                       keywords: Optional[str] = None,
                       comments: Optional[str] = None) -> Dict[str, Any]:
        """
        设置PowerPoint文档属性

        Args:
            path: PPT文件路径
            title: 文档标题
            author: 作者
            subject: 主题
            keywords: 关键词
            comments: 备注

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)
            props = prs.core_properties

            updated = []
            if title is not None:
                props.title = title
                updated.append("title")
            if author is not None:
                props.author = author
                updated.append("author")
            if subject is not None:
                props.subject = subject
                updated.append("subject")
            if keywords is not None:
                props.keywords = keywords
                updated.append("keywords")
            if comments is not None:
                props.comments = comments
                updated.append("comments")

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "updated_properties": updated,
                "message": f"已更新 {len(updated)} 个属性"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("ppt_merge",
                   required_params=['paths', 'output_path'],
                   optional_params=[])
    @staticmethod
    def merge(paths: List[str],
              output_path: str) -> Dict[str, Any]:
        """
        合并多个PowerPoint文档

        Args:
            paths: 要合并的PPT文件路径列表
            output_path: 输出文件路径

        Returns:
            {success, total_slides, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation

            if len(paths) < 2:
                return {"success": False, "error": "至少需要2个文件进行合并"}

            # 检查所有文件存在
            for p in paths:
                if not os.path.exists(p):
                    return {"success": False, "error": f"文件不存在: {p}"}

            # 以第一个文件为基础
            merged_prs = Presentation(paths[0])
            slides_count = [len(merged_prs.slides)]

            # 合并其他文件的幻灯片
            for ppt_path in paths[1:]:
                src_prs = Presentation(ppt_path)
                slides_count.append(len(src_prs.slides))

                for slide in src_prs.slides:
                    # 复制幻灯片布局
                    slide_layout = merged_prs.slide_layouts[6]  # 空白布局

                    # 尝试找到匹配的布局
                    for layout in merged_prs.slide_layouts:
                        if layout.name == slide.slide_layout.name:
                            slide_layout = layout
                            break

                    new_slide = merged_prs.slides.add_slide(slide_layout)

                    # 复制形状
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            # 复制文本框
                            new_shape = new_slide.shapes.add_textbox(
                                shape.left, shape.top, shape.width, shape.height
                            )
                            for i, para in enumerate(shape.text_frame.paragraphs):
                                if i == 0:
                                    new_para = new_shape.text_frame.paragraphs[0]
                                else:
                                    new_para = new_shape.text_frame.add_paragraph()
                                new_para.text = para.text

            # 创建输出目录
            out_dir = os.path.dirname(output_path)
            if out_dir:
                os.makedirs(out_dir, exist_ok=True)
            merged_prs.save(output_path)

            total = len(merged_prs.slides)

            return {
                "success": True,
                "output_path": output_path,
                "merged_files": len(paths),
                "slides_per_file": slides_count,
                "total_slides": total,
                "message": f"已合并 {len(paths)} 个文件，共 {total} 张幻灯片"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    # ========== 幻灯片操作 ==========

    @register_tool("slide_add",
                   required_params=['path'],
                   optional_params=['layout', 'index'])
    @staticmethod
    def slide_add(path: str,
                  layout: Optional[str] = None,
                  index: Optional[int] = None) -> Dict[str, Any]:
        """
        添加新幻灯片

        Args:
            path: PPT文件路径
            layout: 布局名称，如 'Title Slide', 'Title and Content', 'Blank' 等
            index: 插入位置（从1开始），None表示添加到末尾

        Returns:
            {success, slide_index, layout, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            # 查找布局
            slide_layout = None
            layout_name = layout or "Blank"

            for sl in prs.slide_layouts:
                if sl.name.lower() == layout_name.lower():
                    slide_layout = sl
                    break

            # 如果没找到，使用第一个布局或空白布局
            if slide_layout is None:
                # 尝试找空白布局
                for sl in prs.slide_layouts:
                    if 'blank' in sl.name.lower():
                        slide_layout = sl
                        break
                if slide_layout is None:
                    slide_layout = prs.slide_layouts[0]

            # 添加幻灯片
            new_slide = prs.slides.add_slide(slide_layout)
            new_index = len(prs.slides)

            # 如果指定了位置，需要移动幻灯片
            # 注意：python-pptx 不直接支持插入到指定位置
            # 幻灯片会添加到末尾

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide_index": new_index,
                "layout": slide_layout.name,
                "total_slides": len(prs.slides),
                "message": f"已添加幻灯片到第 {new_index} 页，布局: {slide_layout.name}"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("slide_delete",
                   required_params=['path', 'index'],
                   optional_params=[])
    @staticmethod
    def slide_delete(path: str,
                     index: int) -> Dict[str, Any]:
        """
        删除指定幻灯片

        Args:
            path: PPT文件路径
            index: 幻灯片索引（从1开始）

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if index < 1 or index > len(prs.slides):
                return {"success": False, "error": f"索引 {index} 超出范围 (1-{len(prs.slides)})"}

            # 删除幻灯片
            slide_id = prs.slides._sldIdLst[index - 1].rId
            prs.part.drop_rel(slide_id)
            del prs.slides._sldIdLst[index - 1]

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "deleted_index": index,
                "remaining_slides": len(prs.slides),
                "message": f"已删除第 {index} 张幻灯片"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("slide_duplicate",
                   required_params=['path', 'index'],
                   optional_params=[])
    @staticmethod
    def slide_duplicate(path: str,
                        index: int) -> Dict[str, Any]:
        """
        复制指定幻灯片

        Args:
            path: PPT文件路径
            index: 要复制的幻灯片索引（从1开始）

        Returns:
            {success, new_index, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from copy import deepcopy
            import copy

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if index < 1 or index > len(prs.slides):
                return {"success": False, "error": f"索引 {index} 超出范围 (1-{len(prs.slides)})"}

            source_slide = prs.slides[index - 1]

            # 添加新幻灯片使用相同布局
            new_slide = prs.slides.add_slide(source_slide.slide_layout)

            # 复制形状
            for shape in source_slide.shapes:
                if shape.has_text_frame:
                    new_shape = new_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    for i, para in enumerate(shape.text_frame.paragraphs):
                        if i == 0:
                            new_para = new_shape.text_frame.paragraphs[0]
                        else:
                            new_para = new_shape.text_frame.add_paragraph()
                        new_para.text = para.text

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "source_index": index,
                "new_index": len(prs.slides),
                "total_slides": len(prs.slides),
                "message": f"已复制第 {index} 张幻灯片到第 {len(prs.slides)} 页"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("slide_get_layouts",
                   required_params=['path'],
                   optional_params=[])
    @staticmethod
    def slide_get_layouts(path: str) -> Dict[str, Any]:
        """
        获取PPT中可用的幻灯片布局列表

        Args:
            path: PPT文件路径

        Returns:
            {success, layouts, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            layouts = []
            for idx, layout in enumerate(prs.slide_layouts):
                layouts.append({
                    "index": idx,
                    "name": layout.name
                })

            return {
                "success": True,
                "path": path,
                "layouts": layouts,
                "total_layouts": len(layouts),
                "message": f"找到 {len(layouts)} 个可用布局"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    # ========== 内容添加 ==========

    @register_tool("shape_add_text",
                   required_params=['path', 'slide', 'text'],
                   optional_params=['left', 'top', 'width', 'height', 'font_name', 'font_size', 'bold', 'italic', 'color', 'alignment'])
    @staticmethod
    def shape_add_text(path: str,
                       slide: int,
                       text: str,
                       left: str = "1in",
                       top: str = "1in",
                       width: str = "8in",
                       height: str = "1in",
                       font_name: Optional[str] = None,
                       font_size: Optional[int] = None,
                       bold: bool = False,
                       italic: bool = False,
                       color: Optional[str] = None,
                       alignment: Optional[str] = None) -> Dict[str, Any]:
        """
        在幻灯片中添加文本框

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            text: 文本内容
            left: 左边距，如 '1in', '2.54cm'
            top: 上边距
            width: 文本框宽度
            height: 文本框高度
            font_name: 字体名称
            font_size: 字号（磅）
            bold: 是否加粗
            italic: 是否斜体
            color: 字体颜色（十六进制，如 'FF0000'）
            alignment: 对齐方式 (left/center/right)

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引 {slide} 超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            # 添加文本框
            textbox = target_slide.shapes.add_textbox(
                _parse_length(left),
                _parse_length(top),
                _parse_length(width),
                _parse_length(height)
            )

            tf = textbox.text_frame
            tf.word_wrap = True

            # 设置文本
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text

            # 设置字体属性
            font = run.font
            if font_name:
                font.name = font_name
            if font_size:
                font.size = Pt(font_size)
            font.bold = bold
            font.italic = italic

            if color:
                try:
                    color = color.lstrip('#')
                    font.color.rgb = RGBColor(
                        int(color[0:2], 16),
                        int(color[2:4], 16),
                        int(color[4:6], 16)
                    )
                except (ValueError, IndexError):
                    pass

            # 设置对齐
            if alignment:
                align_map = {
                    'left': PP_ALIGN.LEFT,
                    'center': PP_ALIGN.CENTER,
                    'right': PP_ALIGN.RIGHT
                }
                p.alignment = align_map.get(alignment.lower(), PP_ALIGN.LEFT)

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "text": text[:50] + "..." if len(text) > 50 else text,
                "message": f"已在第 {slide} 张幻灯片添加文本框"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("shape_add_image",
                   required_params=['path', 'slide', 'image_path'],
                   optional_params=['left', 'top', 'width', 'height'])
    @staticmethod
    def shape_add_image(path: str,
                        slide: int,
                        image_path: str,
                        left: str = "1in",
                        top: str = "1in",
                        width: Optional[str] = None,
                        height: Optional[str] = None) -> Dict[str, Any]:
        """
        在幻灯片中添加图片

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            image_path: 图片文件路径
            left: 左边距
            top: 上边距
            width: 图片宽度（可选，保持比例）
            height: 图片高度（可选）

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}
            if not os.path.exists(image_path):
                return {"success": False, "error": f"图片文件不存在: {image_path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引 {slide} 超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            # 添加图片
            left_val = _parse_length(left)
            top_val = _parse_length(top)
            width_val = _parse_length(width) if width else None
            height_val = _parse_length(height) if height else None

            picture = target_slide.shapes.add_picture(
                image_path,
                left_val,
                top_val,
                width_val,
                height_val
            )

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "image_path": image_path,
                "position": {"left": left, "top": top},
                "message": f"已在第 {slide} 张幻灯片添加图片"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("shape_add_table",
                   required_params=['path', 'slide', 'rows', 'cols'],
                   optional_params=['left', 'top', 'width', 'height', 'data'])
    @staticmethod
    def shape_add_table(path: str,
                        slide: int,
                        rows: int,
                        cols: int,
                        left: str = "1in",
                        top: str = "2in",
                        width: str = "8in",
                        height: str = "3in",
                        data: Optional[List[List[str]]] = None) -> Dict[str, Any]:
        """
        在幻灯片中添加表格

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            rows: 行数
            cols: 列数
            left: 左边距
            top: 上边距
            width: 表格宽度
            height: 表格高度
            data: 表格数据（二维数组）

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引 {slide} 超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            # 添加表格
            table_shape = target_slide.shapes.add_table(
                rows, cols,
                _parse_length(left),
                _parse_length(top),
                _parse_length(width),
                _parse_length(height)
            )

            table = table_shape.table

            # 填充数据
            if data:
                for r_idx, row_data in enumerate(data):
                    if r_idx >= rows:
                        break
                    for c_idx, cell_text in enumerate(row_data):
                        if c_idx >= cols:
                            break
                        table.cell(r_idx, c_idx).text = str(cell_text) if cell_text else ""

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "table_size": {"rows": rows, "cols": cols},
                "message": f"已在第 {slide} 张幻灯片添加 {rows}x{cols} 表格"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("shape_add_shape",
                   required_params=['path', 'slide', 'shape_type'],
                   optional_params=['left', 'top', 'width', 'height', 'fill_color', 'line_color', 'text'])
    @staticmethod
    def shape_add_shape(path: str,
                        slide: int,
                        shape_type: str,
                        left: str = "2in",
                        top: str = "2in",
                        width: str = "2in",
                        height: str = "2in",
                        fill_color: Optional[str] = None,
                        line_color: Optional[str] = None,
                        text: Optional[str] = None) -> Dict[str, Any]:
        """
        在幻灯片中添加形状

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            shape_type: 形状类型 (rectangle/oval/triangle/arrow_right/star等)
            left: 左边距
            top: 上边距
            width: 形状宽度
            height: 形状高度
            fill_color: 填充颜色（十六进制）
            line_color: 边框颜色（十六进制）
            text: 形状内的文字

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引 {slide} 超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            # 形状类型映射
            shape_map = {
                'rectangle': MSO_SHAPE.RECTANGLE,
                'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
                'oval': MSO_SHAPE.OVAL,
                'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
                'right_triangle': MSO_SHAPE.RIGHT_TRIANGLE,
                'diamond': MSO_SHAPE.DIAMOND,
                'pentagon': MSO_SHAPE.PENTAGON,
                'hexagon': MSO_SHAPE.HEXAGON,
                'arrow_right': MSO_SHAPE.RIGHT_ARROW,
                'arrow_left': MSO_SHAPE.LEFT_ARROW,
                'arrow_up': MSO_SHAPE.UP_ARROW,
                'arrow_down': MSO_SHAPE.DOWN_ARROW,
                'star': MSO_SHAPE.STAR_5_POINT,
                'star_4': MSO_SHAPE.STAR_4_POINT,
                'star_6': MSO_SHAPE.STAR_6_POINT,
                'heart': MSO_SHAPE.HEART,
                'lightning': MSO_SHAPE.LIGHTNING_BOLT,
                'cloud': MSO_SHAPE.CLOUD,
            }

            mso_shape = shape_map.get(shape_type.lower(), MSO_SHAPE.RECTANGLE)

            # 添加形状
            shape = target_slide.shapes.add_shape(
                mso_shape,
                _parse_length(left),
                _parse_length(top),
                _parse_length(width),
                _parse_length(height)
            )

            # 设置填充颜色
            if fill_color:
                try:
                    fill_color = fill_color.lstrip('#')
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(
                        int(fill_color[0:2], 16),
                        int(fill_color[2:4], 16),
                        int(fill_color[4:6], 16)
                    )
                except (ValueError, IndexError):
                    pass

            # 设置边框颜色
            if line_color:
                try:
                    line_color = line_color.lstrip('#')
                    shape.line.color.rgb = RGBColor(
                        int(line_color[0:2], 16),
                        int(line_color[2:4], 16),
                        int(line_color[4:6], 16)
                    )
                except (ValueError, IndexError):
                    pass

            # 添加文字
            if text and shape.has_text_frame:
                shape.text = text

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "shape_type": shape_type,
                "message": f"已在第 {slide} 张幻灯片添加{shape_type}形状"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    # ========== 高级功能 ==========

    @register_tool("slide_set_background",
                   required_params=['path', 'slide'],
                   optional_params=['color', 'image_path'])
    @staticmethod
    def slide_set_background(path: str,
                             slide: int,
                             color: Optional[str] = None,
                             image_path: Optional[str] = None) -> Dict[str, Any]:
        """
        设置幻灯片背景

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            color: 背景颜色（十六进制，如 'FFFFFF'）
            image_path: 背景图片路径

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from pptx.dml.color import RGBColor

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引 {slide} 超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]
            background = target_slide.background

            if color:
                color = color.lstrip('#')
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(
                    int(color[0:2], 16),
                    int(color[2:4], 16),
                    int(color[4:6], 16)
                )
                prs.save(path)
                return {
                    "success": True,
                    "path": path,
                    "slide": slide,
                    "background_type": "color",
                    "color": color,
                    "message": f"已设置第 {slide} 张幻灯片背景颜色"
                }

            # 图片背景需要更复杂的处理
            if image_path:
                return {
                    "success": False,
                    "error": "图片背景功能暂不支持，请使用颜色背景或手动设置"
                }

            return {"success": False, "error": "请指定 color 或 image_path"}

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("slide_add_notes",
                   required_params=['path', 'slide', 'notes'],
                   optional_params=[])
    @staticmethod
    def slide_add_notes(path: str,
                        slide: int,
                        notes: str) -> Dict[str, Any]:
        """
        为幻灯片添加演讲者备注

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            notes: 备注内容

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引 {slide} 超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            # 获取或创建备注页
            notes_slide = target_slide.notes_slide
            notes_slide.notes_text_frame.text = notes

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "notes_length": len(notes),
                "message": f"已为第 {slide} 张幻灯片添加备注"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("ppt_status",
                   required_params=[],
                   optional_params=[])
    @staticmethod
    def get_status() -> Dict[str, Any]:
        """
        获取PPT模块状态

        Returns:
            {success, pptx_available, version, features, message}
        """
        try:
            pptx_available = check_import("pptx")
            version = None

            if pptx_available:
                try:
                    import pptx
                    version = pptx.__version__
                except (AttributeError, ImportError):
                    version = "unknown"

            features = []
            if pptx_available:
                features.extend([
                    "PPT创建与读取",
                    "幻灯片管理（添加/删除/复制）",
                    "文本框添加",
                    "图片添加",
                    "表格添加",
                    "形状添加",
                    "背景设置",
                    "演讲者备注",
                    "PPT合并"
                ])

            return {
                "success": True,
                "pptx_available": pptx_available,
                "version": version,
                "features": features,
                "message": f"python-pptx {'已安装 v' + version if pptx_available else '未安装'}"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    # ========== 母版操作 ==========

    @register_tool("master_list",
                   required_params=['path'],
                   optional_params=[])
    @staticmethod
    def master_list(path: str) -> Dict[str, Any]:
        """
        列出所有母版和布局

        Args:
            path: PPT文件路径

        Returns:
            {success, masters: [{name, layouts}], message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)
            masters = []

            for idx, master in enumerate(prs.slide_masters):
                layouts = []
                for layout_idx, layout in enumerate(master.slide_layouts):
                    placeholders = []
                    for ph in layout.placeholders:
                        placeholders.append({
                            "idx": ph.placeholder_format.idx,
                            "type": str(ph.placeholder_format.type).split('.')[-1].strip('()'),
                            "name": ph.name
                        })
                    layouts.append({
                        "index": layout_idx,
                        "name": layout.name,
                        "placeholder_count": len(list(layout.placeholders))
                    })

                masters.append({
                    "index": idx,
                    "name": master.name if hasattr(master, 'name') else f"母版 {idx + 1}",
                    "layout_count": len(layouts),
                    "layouts": layouts
                })

            return {
                "success": True,
                "path": path,
                "master_count": len(masters),
                "masters": masters,
                "message": f"找到 {len(masters)} 个母版"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("master_get_info",
                   required_params=['path'],
                   optional_params=['master_index'])
    @staticmethod
    def master_get_info(path: str, master_index: int = 0) -> Dict[str, Any]:
        """
        获取母版详细信息

        Args:
            path: PPT文件路径
            master_index: 母版索引（默认0）

        Returns:
            {success, master_info, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from pptx.util import Emu

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if master_index < 0 or master_index >= len(prs.slide_masters):
                return {"success": False, "error": f"母版索引超出范围 (0-{len(prs.slide_masters)-1})"}

            master = prs.slide_masters[master_index]

            # 获取母版中的形状
            shapes_info = []
            for shape in master.shapes:
                shape_info = {
                    "shape_id": shape.shape_id,
                    "name": shape.name,
                    "type": shape.shape_type.name if hasattr(shape.shape_type, 'name') else str(shape.shape_type),
                    "left": round(shape.left / Emu(914400), 2),  # 转换为英寸
                    "top": round(shape.top / Emu(914400), 2),
                    "width": round(shape.width / Emu(914400), 2),
                    "height": round(shape.height / Emu(914400), 2)
                }
                if shape.has_text_frame:
                    shape_info["text"] = shape.text_frame.text[:50] if shape.text_frame.text else ""
                shapes_info.append(shape_info)

            # 获取母版占位符
            placeholders = []
            for ph in master.placeholders:
                placeholders.append({
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type).split('.')[-1].strip('()'),
                    "name": ph.name,
                    "has_text": ph.has_text_frame
                })

            # 获取布局信息
            layouts = []
            for idx, layout in enumerate(master.slide_layouts):
                layouts.append({
                    "index": idx,
                    "name": layout.name
                })

            return {
                "success": True,
                "path": path,
                "master_index": master_index,
                "shape_count": len(shapes_info),
                "shapes": shapes_info,
                "placeholder_count": len(placeholders),
                "placeholders": placeholders,
                "layout_count": len(layouts),
                "layouts": layouts,
                "message": f"母版 {master_index} 包含 {len(shapes_info)} 个形状，{len(layouts)} 个布局"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("placeholder_list",
                   required_params=['path', 'slide'],
                   optional_params=[])
    @staticmethod
    def placeholder_list(path: str, slide: int) -> Dict[str, Any]:
        """
        列出幻灯片中的占位符

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）

        Returns:
            {success, placeholders, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from pptx.util import Emu

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]
            placeholders = []

            for ph in target_slide.placeholders:
                ph_info = {
                    "idx": ph.placeholder_format.idx,
                    "type": str(ph.placeholder_format.type).split('.')[-1].strip('()'),
                    "name": ph.name,
                    "left": round(ph.left / Emu(914400), 2),
                    "top": round(ph.top / Emu(914400), 2),
                    "width": round(ph.width / Emu(914400), 2),
                    "height": round(ph.height / Emu(914400), 2),
                    "has_text": ph.has_text_frame
                }
                if ph.has_text_frame:
                    ph_info["text"] = ph.text_frame.text[:100] if ph.text_frame.text else ""
                placeholders.append(ph_info)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "layout_name": target_slide.slide_layout.name,
                "placeholder_count": len(placeholders),
                "placeholders": placeholders,
                "message": f"第 {slide} 张幻灯片 (布局: {target_slide.slide_layout.name}) 包含 {len(placeholders)} 个占位符"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("placeholder_set",
                   required_params=['path', 'slide', 'idx'],
                   optional_params=['text', 'font_name', 'font_size', 'bold', 'italic', 'color'])
    @staticmethod
    def placeholder_set(path: str,
                        slide: int,
                        idx: int,
                        text: Optional[str] = None,
                        font_name: Optional[str] = None,
                        font_size: Optional[int] = None,
                        bold: bool = False,
                        italic: bool = False,
                        color: Optional[str] = None) -> Dict[str, Any]:
        """
        设置占位符内容

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            idx: 占位符索引
            text: 文本内容
            font_name: 字体名称
            font_size: 字号
            bold: 是否加粗
            italic: 是否斜体
            color: 字体颜色（十六进制）

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from pptx.util import Pt
            from pptx.dml.color import RGBColor

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            # 查找占位符
            placeholder = None
            for ph in target_slide.placeholders:
                if ph.placeholder_format.idx == idx:
                    placeholder = ph
                    break

            if placeholder is None:
                return {"success": False, "error": f"未找到索引为 {idx} 的占位符"}

            if not placeholder.has_text_frame:
                return {"success": False, "error": f"占位符 {idx} 不支持文本"}

            # 设置文本
            if text is not None:
                placeholder.text_frame.clear()
                p = placeholder.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text

                # 设置字体
                if font_name:
                    run.font.name = font_name
                if font_size:
                    run.font.size = Pt(font_size)
                if bold:
                    run.font.bold = True
                if italic:
                    run.font.italic = True
                if color:
                    color = color.lstrip('#')
                    run.font.color.rgb = RGBColor(
                        int(color[0:2], 16),
                        int(color[2:4], 16),
                        int(color[4:6], 16)
                    )

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "placeholder_idx": idx,
                "message": f"已设置第 {slide} 张幻灯片的占位符 {idx}"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    # ========== 动画操作 ==========

    @register_tool("animation_add",
                   required_params=['path', 'slide', 'shape_index', 'effect'],
                   optional_params=['trigger', 'duration', 'delay', 'direction'])
    @staticmethod
    def animation_add(path: str,
                      slide: int,
                      shape_index: int,
                      effect: str,
                      trigger: str = "on_click",
                      duration: float = 0.5,
                      delay: float = 0.0,
                      direction: Optional[str] = None) -> Dict[str, Any]:
        """
        为形状添加动画效果

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            shape_index: 形状索引（从0开始，可通过placeholder_list获取）
            effect: 动画效果 (appear/fade/fly_in/float_in/zoom/wipe/fade_out/fly_out等)
            trigger: 触发方式 (on_click/with_previous/after_previous)
            duration: 持续时间（秒）
            delay: 延迟时间（秒）
            direction: 方向 (left/right/top/bottom，仅fly_in/fly_out等支持)

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from lxml import etree

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]
            shapes = list(target_slide.shapes)

            if shape_index < 0 or shape_index >= len(shapes):
                return {"success": False, "error": f"形状索引超出范围 (0-{len(shapes)-1})"}

            target_shape = shapes[shape_index]
            shape_id = target_shape.shape_id

            # 动画效果映射
            effect_map = {
                # 进入动画
                'appear': ('entr', '1', None),
                'fade': ('entr', '10', None),
                'fly_in': ('entr', '2', 'from'),
                'float_in': ('entr', '42', 'from'),
                'zoom': ('entr', '23', None),
                'wipe': ('entr', '22', 'from'),
                'split': ('entr', '16', None),
                'wheel': ('entr', '21', None),
                # 强调动画
                'pulse': ('emph', '26', None),
                'spin': ('emph', '8', None),
                'grow_shrink': ('emph', '6', None),
                # 退出动画
                'disappear': ('exit', '1', None),
                'fade_out': ('exit', '10', None),
                'fly_out': ('exit', '2', 'to'),
                'zoom_out': ('exit', '23', None),
            }

            if effect.lower() not in effect_map:
                available = ', '.join(effect_map.keys())
                return {"success": False, "error": f"不支持的动画效果: {effect}。可用: {available}"}

            anim_type, preset_id, dir_attr = effect_map[effect.lower()]

            # 方向映射
            direction_map = {
                'left': 'l', 'right': 'r', 'top': 'u', 'up': 'u',
                'bottom': 'd', 'down': 'd',
                'top_left': 'lu', 'top_right': 'ru',
                'bottom_left': 'ld', 'bottom_right': 'rd'
            }
            dir_value = direction_map.get(direction.lower() if direction else 'bottom', 'd')

            # 触发方式映射
            trigger_map = {
                'on_click': 'indefinite',
                'with_previous': '0',
                'after_previous': '0'
            }

            # PowerPoint XML命名空间
            nsmap = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
            }

            # 获取幻灯片XML元素
            slide_elem = target_slide._element

            # 查找或创建timing元素
            timing = slide_elem.find('.//p:timing', nsmap)
            if timing is None:
                timing = etree.SubElement(slide_elem, '{%s}timing' % nsmap['p'])

            # 查找或创建tnLst
            tnLst = timing.find('p:tnLst', nsmap)
            if tnLst is None:
                tnLst = etree.SubElement(timing, '{%s}tnLst' % nsmap['p'])

            # 查找或创建par（根时间节点）
            par = tnLst.find('p:par', nsmap)
            if par is None:
                par = etree.SubElement(tnLst, '{%s}par' % nsmap['p'])
                cTn_root = etree.SubElement(par, '{%s}cTn' % nsmap['p'],
                                            id="1", dur="indefinite", restart="never", nodeType="tmRoot")
                childTnLst_root = etree.SubElement(cTn_root, '{%s}childTnLst' % nsmap['p'])
                seq = etree.SubElement(childTnLst_root, '{%s}seq' % nsmap['p'], concurrent="1", nextAc="seek")
                cTn_seq = etree.SubElement(seq, '{%s}cTn' % nsmap['p'], id="2", dur="indefinite", nodeType="mainSeq")
                main_childTnLst = etree.SubElement(cTn_seq, '{%s}childTnLst' % nsmap['p'])
            else:
                # 找到mainSeq的childTnLst
                main_childTnLst = par.find('.//p:cTn[@nodeType="mainSeq"]/p:childTnLst', nsmap)
                if main_childTnLst is None:
                    cTn_root = par.find('p:cTn', nsmap)
                    if cTn_root is None:
                        cTn_root = etree.SubElement(par, '{%s}cTn' % nsmap['p'],
                                                    id="1", dur="indefinite", restart="never", nodeType="tmRoot")
                    childTnLst_root = cTn_root.find('p:childTnLst', nsmap)
                    if childTnLst_root is None:
                        childTnLst_root = etree.SubElement(cTn_root, '{%s}childTnLst' % nsmap['p'])
                    seq = childTnLst_root.find('p:seq', nsmap)
                    if seq is None:
                        seq = etree.SubElement(childTnLst_root, '{%s}seq' % nsmap['p'], concurrent="1", nextAc="seek")
                    cTn_seq = seq.find('p:cTn', nsmap)
                    if cTn_seq is None:
                        cTn_seq = etree.SubElement(seq, '{%s}cTn' % nsmap['p'], id="2", dur="indefinite", nodeType="mainSeq")
                    main_childTnLst = cTn_seq.find('p:childTnLst', nsmap)
                    if main_childTnLst is None:
                        main_childTnLst = etree.SubElement(cTn_seq, '{%s}childTnLst' % nsmap['p'])

            # 计算新的ID
            all_ids = [int(elem.get('id', 0)) for elem in slide_elem.iter() if elem.get('id', '').isdigit()]
            new_id = max(all_ids) + 1 if all_ids else 3

            # 转换时间为毫秒
            dur_ms = str(int(duration * 1000))
            delay_ms = str(int(delay * 1000))

            # 创建动画序列
            par_anim = etree.SubElement(main_childTnLst, '{%s}par' % nsmap['p'])
            cTn_anim = etree.SubElement(par_anim, '{%s}cTn' % nsmap['p'], id=str(new_id), fill="hold")

            # 设置触发条件
            stCondLst = etree.SubElement(cTn_anim, '{%s}stCondLst' % nsmap['p'])
            if trigger == 'on_click':
                etree.SubElement(stCondLst, '{%s}cond' % nsmap['p'], delay="indefinite")
            elif trigger == 'after_previous':
                etree.SubElement(stCondLst, '{%s}cond' % nsmap['p'], evt="onEnd", delay=delay_ms)
            else:  # with_previous
                etree.SubElement(stCondLst, '{%s}cond' % nsmap['p'], delay=delay_ms)

            # 创建子动画节点
            childTnLst_anim = etree.SubElement(cTn_anim, '{%s}childTnLst' % nsmap['p'])
            par_effect = etree.SubElement(childTnLst_anim, '{%s}par' % nsmap['p'])
            cTn_effect = etree.SubElement(par_effect, '{%s}cTn' % nsmap['p'],
                                          id=str(new_id + 1), presetID=preset_id,
                                          presetClass=anim_type, presetSubtype="1",
                                          fill="hold", nodeType="clickEffect")

            if dir_attr and direction:
                cTn_effect.set('presetSubtype', {'l': '8', 'r': '2', 'u': '4', 'd': '1'}.get(dir_value, '1'))

            stCondLst2 = etree.SubElement(cTn_effect, '{%s}stCondLst' % nsmap['p'])
            etree.SubElement(stCondLst2, '{%s}cond' % nsmap['p'], delay="0")

            childTnLst_effect = etree.SubElement(cTn_effect, '{%s}childTnLst' % nsmap['p'])

            # 添加set动画（visibility）
            set_elem = etree.SubElement(childTnLst_effect, '{%s}set' % nsmap['p'])
            cBhvr_set = etree.SubElement(set_elem, '{%s}cBhvr' % nsmap['p'])
            cTn_set = etree.SubElement(cBhvr_set, '{%s}cTn' % nsmap['p'], id=str(new_id + 2), dur="1", fill="hold")
            etree.SubElement(cTn_set, '{%s}stCondLst' % nsmap['p']).append(
                etree.Element('{%s}cond' % nsmap['p'], delay="0"))
            tgtEl_set = etree.SubElement(cBhvr_set, '{%s}tgtEl' % nsmap['p'])
            etree.SubElement(tgtEl_set, '{%s}spTgt' % nsmap['p'], spid=str(shape_id))
            attrNameLst = etree.SubElement(cBhvr_set, '{%s}attrNameLst' % nsmap['p'])
            etree.SubElement(attrNameLst, '{%s}attrName' % nsmap['p']).text = "style.visibility"
            to_elem = etree.SubElement(set_elem, '{%s}to' % nsmap['p'])
            etree.SubElement(to_elem, '{%s}strVal' % nsmap['p'], val="visible")

            # 添加anim动画（主效果）
            if effect.lower() in ['fade', 'fade_out']:
                # 淡入淡出使用animEffect
                animEffect = etree.SubElement(childTnLst_effect, '{%s}animEffect' % nsmap['p'],
                                              transition="in" if anim_type == 'entr' else "out",
                                              filter="fade")
                cBhvr_anim = etree.SubElement(animEffect, '{%s}cBhvr' % nsmap['p'])
                cTn_anim2 = etree.SubElement(cBhvr_anim, '{%s}cTn' % nsmap['p'],
                                             id=str(new_id + 3), dur=dur_ms)
                tgtEl_anim = etree.SubElement(cBhvr_anim, '{%s}tgtEl' % nsmap['p'])
                etree.SubElement(tgtEl_anim, '{%s}spTgt' % nsmap['p'], spid=str(shape_id))

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "shape_index": shape_index,
                "shape_id": shape_id,
                "effect": effect,
                "trigger": trigger,
                "duration": duration,
                "message": f"已为第 {slide} 张幻灯片的形状 {shape_index} 添加 {effect} 动画"
            }

        except ImportError:
            return {"success": False, "error": "需要安装lxml: pip install lxml"}
        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("animation_list",
                   required_params=['path', 'slide'],
                   optional_params=[])
    @staticmethod
    def animation_list(path: str, slide: int) -> Dict[str, Any]:
        """
        列出幻灯片上的动画

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）

        Returns:
            {success, animations, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from lxml import etree

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            # 获取形状ID到名称的映射
            shape_map = {shape.shape_id: shape.name for shape in target_slide.shapes}

            # PowerPoint XML命名空间
            nsmap = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
            }

            slide_elem = target_slide._element
            animations = []

            # 动画类型映射
            preset_class_map = {
                'entr': '进入',
                'emph': '强调',
                'exit': '退出',
                'path': '路径'
            }

            preset_id_map = {
                '1': 'appear/disappear',
                '2': 'fly',
                '10': 'fade',
                '16': 'split',
                '21': 'wheel',
                '22': 'wipe',
                '23': 'zoom',
                '26': 'pulse',
                '42': 'float',
            }

            # 查找所有动画节点
            for cTn in slide_elem.iter('{%s}cTn' % nsmap['p']):
                preset_id = cTn.get('presetID')
                preset_class = cTn.get('presetClass')

                if preset_id and preset_class:
                    # 查找目标形状
                    spTgt = cTn.find('.//p:spTgt', nsmap)
                    if spTgt is not None:
                        shape_id = int(spTgt.get('spid', 0))
                        shape_name = shape_map.get(shape_id, f"Shape {shape_id}")

                        anim_info = {
                            "id": cTn.get('id'),
                            "shape_id": shape_id,
                            "shape_name": shape_name,
                            "type": preset_class_map.get(preset_class, preset_class),
                            "effect": preset_id_map.get(preset_id, f"preset_{preset_id}"),
                            "duration": cTn.get('dur', 'unknown')
                        }
                        animations.append(anim_info)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "animation_count": len(animations),
                "animations": animations,
                "message": f"第 {slide} 张幻灯片有 {len(animations)} 个动画"
            }

        except ImportError:
            return {"success": False, "error": "需要安装lxml: pip install lxml"}
        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("animation_remove",
                   required_params=['path', 'slide'],
                   optional_params=['shape_index', 'remove_all'])
    @staticmethod
    def animation_remove(path: str,
                         slide: int,
                         shape_index: Optional[int] = None,
                         remove_all: bool = False) -> Dict[str, Any]:
        """
        删除动画

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            shape_index: 形状索引（删除该形状的动画）
            remove_all: 是否删除所有动画

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from lxml import etree

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            nsmap = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
            }

            slide_elem = target_slide._element

            if remove_all:
                # 删除整个timing元素
                timing = slide_elem.find('p:timing', nsmap)
                if timing is not None:
                    slide_elem.remove(timing)
                    prs.save(path)
                    return {
                        "success": True,
                        "path": path,
                        "slide": slide,
                        "message": f"已删除第 {slide} 张幻灯片的所有动画"
                    }
                else:
                    return {
                        "success": True,
                        "path": path,
                        "slide": slide,
                        "message": f"第 {slide} 张幻灯片没有动画"
                    }

            if shape_index is not None:
                shapes = list(target_slide.shapes)
                if shape_index < 0 or shape_index >= len(shapes):
                    return {"success": False, "error": f"形状索引超出范围 (0-{len(shapes)-1})"}

                target_shape_id = str(shapes[shape_index].shape_id)
                removed_count = 0

                # 查找并删除与该形状相关的动画
                for spTgt in list(slide_elem.iter('{%s}spTgt' % nsmap['p'])):
                    if spTgt.get('spid') == target_shape_id:
                        # 向上查找par元素并删除
                        parent = spTgt.getparent()
                        while parent is not None:
                            if parent.tag == '{%s}par' % nsmap['p']:
                                grandparent = parent.getparent()
                                if grandparent is not None:
                                    grandparent.remove(parent)
                                    removed_count += 1
                                break
                            parent = parent.getparent()

                prs.save(path)
                return {
                    "success": True,
                    "path": path,
                    "slide": slide,
                    "shape_index": shape_index,
                    "removed_count": removed_count,
                    "message": f"已删除形状 {shape_index} 的 {removed_count} 个动画"
                }

            return {"success": False, "error": "请指定 shape_index 或设置 remove_all=True"}

        except ImportError:
            return {"success": False, "error": "需要安装lxml: pip install lxml"}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ========== 幻灯片切换效果 ==========

    @register_tool("slide_set_transition",
                   required_params=['path', 'slide', 'effect'],
                   optional_params=['speed', 'advance_click', 'advance_time', 'duration'])
    @staticmethod
    def slide_set_transition(path: str,
                             slide: int,
                             effect: str,
                             speed: str = 'medium',
                             advance_click: bool = True,
                             advance_time: Optional[int] = None,
                             duration: Optional[int] = None) -> Dict[str, Any]:
        """
        设置幻灯片切换效果

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            effect: 切换效果 (fade/push/wipe/split/cut/dissolve/cover/uncover/randomBars/blinds/wheel/comb/checker)
            speed: 速度 (slow/medium/fast)
            advance_click: 是否点击切换
            advance_time: 自动切换时间（毫秒）
            duration: 切换持续时间（毫秒）

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from lxml import etree

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]
            slide_elem = target_slide._element

            # 命名空间
            p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            p_tag = '{%s}' % p_ns

            # 效果映射：效果名 -> 子元素tag名
            effect_map = {
                'fade': 'fade',
                'push': 'push',
                'wipe': 'wipe',
                'split': 'split',
                'cut': 'cut',
                'dissolve': 'dissolve',
                'cover': 'cover',
                'uncover': 'uncover',
                'randombars': 'randomBar',
                'random_bars': 'randomBar',
                'blinds': 'blinds',
                'wheel': 'wheel',
                'comb': 'comb',
                'checker': 'checker',
                'random': 'random',
                'strips': 'strips',
                'newsflash': 'newsflash',
                'plus': 'plus',
                'circle': 'circle',
                'diamond': 'diamond',
                'wedge': 'wedge',
                'zoom': 'zoom',
            }

            effect_lower = effect.lower()
            if effect_lower not in effect_map:
                available = ', '.join(sorted(effect_map.keys()))
                return {"success": False, "error": f"不支持的切换效果: {effect}。可用: {available}"}

            effect_tag = effect_map[effect_lower]

            # 速度映射
            speed_map = {'slow': 'slow', 'medium': 'med', 'fast': 'fast'}
            spd = speed_map.get(speed.lower(), 'med')

            # 移除已有的 <p:transition>
            for old_trans in slide_elem.findall(p_tag + 'transition'):
                slide_elem.remove(old_trans)

            # 创建 <p:transition> 元素
            transition = etree.SubElement(slide_elem, p_tag + 'transition')
            transition.set('spd', spd)

            if advance_click:
                transition.set('advClick', '1')
            else:
                transition.set('advClick', '0')

            if advance_time is not None:
                transition.set('advTm', str(advance_time))

            if duration is not None:
                # duration 以毫秒为单位在 OOXML 中
                transition.set('dur', str(duration))

            # 添加效果子元素
            effect_elem = etree.SubElement(transition, p_tag + effect_tag)

            # <p:transition> 需要在 <p:cSld> 之后、<p:timing> 之前
            # 先移除再插入到正确位置
            slide_elem.remove(transition)

            cSld = slide_elem.find(p_tag + 'cSld')
            if cSld is not None:
                cSld_idx = list(slide_elem).index(cSld)
                slide_elem.insert(cSld_idx + 1, transition)
            else:
                # 如果没有 cSld（不太可能），查找 timing
                timing = slide_elem.find(p_tag + 'timing')
                if timing is not None:
                    timing_idx = list(slide_elem).index(timing)
                    slide_elem.insert(timing_idx, transition)
                else:
                    slide_elem.append(transition)

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "effect": effect,
                "speed": speed,
                "advance_click": advance_click,
                "advance_time": advance_time,
                "message": f"已为第 {slide} 张幻灯片设置 {effect} 切换效果"
            }

        except ImportError:
            return {"success": False, "error": "需要安装lxml: pip install lxml"}
        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("slide_remove_transition",
                   required_params=['path', 'slide'],
                   optional_params=[])
    @staticmethod
    def slide_remove_transition(path: str,
                                slide: int) -> Dict[str, Any]:
        """
        移除幻灯片切换效果

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）

        Returns:
            {success, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from lxml import etree

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]
            slide_elem = target_slide._element

            p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            p_tag = '{%s}' % p_ns

            # 移除所有 <p:transition>
            removed = 0
            for trans in slide_elem.findall(p_tag + 'transition'):
                slide_elem.remove(trans)
                removed += 1

            if removed == 0:
                return {
                    "success": True,
                    "path": path,
                    "slide": slide,
                    "message": f"第 {slide} 张幻灯片没有切换效果"
                }

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "message": f"已移除第 {slide} 张幻灯片的切换效果"
            }

        except ImportError:
            return {"success": False, "error": "需要安装lxml: pip install lxml"}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ========== 图表操作 ==========

    # 图表类型映射
    CHART_TYPE_MAP = {
        'column': 'COLUMN_CLUSTERED',
        'column_stacked': 'COLUMN_STACKED',
        'column_stacked_100': 'COLUMN_STACKED_100',
        'bar': 'BAR_CLUSTERED',
        'bar_stacked': 'BAR_STACKED',
        'bar_stacked_100': 'BAR_STACKED_100',
        'line': 'LINE',
        'line_stacked': 'LINE_STACKED',
        'line_markers': 'LINE_MARKERS',
        'pie': 'PIE',
        'pie_exploded': 'PIE_EXPLODED',
        'doughnut': 'DOUGHNUT',
        'doughnut_exploded': 'DOUGHNUT_EXPLODED',
        'area': 'AREA',
        'area_stacked': 'AREA_STACKED',
        'area_stacked_100': 'AREA_STACKED_100',
        'scatter': 'XY_SCATTER',
        'scatter_lines': 'XY_SCATTER_LINES',
        'scatter_smooth': 'XY_SCATTER_SMOOTH',
        'bubble': 'BUBBLE',
        'radar': 'RADAR',
        'radar_filled': 'RADAR_FILLED',
    }

    # 图例位置映射
    LEGEND_POSITION_MAP = {
        'right': 'RIGHT',
        'left': 'LEFT',
        'top': 'TOP',
        'bottom': 'BOTTOM',
        'corner': 'CORNER',
    }

    # 数据标签位置映射
    LABEL_POSITION_MAP = {
        'center': 'CENTER',
        'inside_end': 'INSIDE_END',
        'inside_base': 'INSIDE_BASE',
        'outside_end': 'OUTSIDE_END',
        'best_fit': 'BEST_FIT',
    }

    @register_tool("chart_add",
                   required_params=['path', 'slide', 'chart_type', 'categories', 'series'],
                   optional_params=['x', 'y', 'width', 'height', 'title', 'has_legend',
                                    'legend_position', 'has_data_labels', 'data_label_position'])
    @staticmethod
    def chart_add(path: str,
                  slide: int,
                  chart_type: str,
                  categories: List[str],
                  series: List[Dict[str, Any]],
                  x: Optional[str] = None,
                  y: Optional[str] = None,
                  width: Optional[str] = None,
                  height: Optional[str] = None,
                  title: Optional[str] = None,
                  has_legend: bool = True,
                  legend_position: Optional[str] = None,
                  has_data_labels: bool = False,
                  data_label_position: Optional[str] = None) -> Dict[str, Any]:
        """
        向幻灯片添加图表

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            chart_type: 图表类型 (column, bar, line, pie, area, scatter, bubble, doughnut, radar等)
            categories: 分类标签列表，如 ['Q1', 'Q2', 'Q3']
            series: 系列数据列表，如 [{"name": "销售额", "values": [100, 200, 150]}]
            x: 图表左边距，如 '1in', '2.5cm'
            y: 图表上边距
            width: 图表宽度
            height: 图表高度
            title: 图表标题
            has_legend: 是否显示图例
            legend_position: 图例位置 (right, left, top, bottom, corner)
            has_data_labels: 是否显示数据标签
            data_label_position: 数据标签位置 (center, inside_end, outside_end等)

        Returns:
            {success, path, slide, chart_type, series_count, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from pptx.util import Inches
            from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
            from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            # 获取图表类型
            chart_type_lower = chart_type.lower()
            if chart_type_lower not in PPTOperations.CHART_TYPE_MAP:
                available_types = ', '.join(PPTOperations.CHART_TYPE_MAP.keys())
                return {"success": False, "error": f"不支持的图表类型: {chart_type}。支持的类型: {available_types}"}

            xl_chart_type_name = PPTOperations.CHART_TYPE_MAP[chart_type_lower]
            xl_chart_type = getattr(XL_CHART_TYPE, xl_chart_type_name)

            # 设置位置和尺寸（默认值）
            chart_x = _parse_length(x) if x else Inches(1)
            chart_y = _parse_length(y) if y else Inches(1.5)
            chart_width = _parse_length(width) if width else Inches(8)
            chart_height = _parse_length(height) if height else Inches(5)

            # 根据图表类型选择数据类
            is_xy_chart = chart_type_lower in ('scatter', 'scatter_lines', 'scatter_smooth')
            is_bubble_chart = chart_type_lower == 'bubble'

            if is_xy_chart:
                # XY散点图
                chart_data = XyChartData()
                for s in series:
                    s_name = s.get('name', 'Series')
                    s_values = s.get('values', [])
                    xy_series = chart_data.add_series(s_name)
                    # values应该是[(x1,y1), (x2,y2), ...] 或者 [[x1,y1], [x2,y2], ...]
                    for point in s_values:
                        if isinstance(point, (list, tuple)) and len(point) >= 2:
                            xy_series.add_data_point(point[0], point[1])

            elif is_bubble_chart:
                # 气泡图
                chart_data = BubbleChartData()
                for s in series:
                    s_name = s.get('name', 'Series')
                    s_values = s.get('values', [])
                    bubble_series = chart_data.add_series(s_name)
                    # values应该是[(x,y,size), ...] 或者 [[x,y,size], ...]
                    for point in s_values:
                        if isinstance(point, (list, tuple)) and len(point) >= 3:
                            bubble_series.add_data_point(point[0], point[1], point[2])

            else:
                # 分类图表（柱状图、折线图、饼图等）
                chart_data = CategoryChartData()
                chart_data.categories = categories
                for s in series:
                    s_name = s.get('name', 'Series')
                    s_values = s.get('values', [])
                    chart_data.add_series(s_name, s_values)

            # 添加图表
            graphic_frame = target_slide.shapes.add_chart(
                xl_chart_type, chart_x, chart_y, chart_width, chart_height, chart_data
            )
            chart = graphic_frame.chart

            # 设置标题
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.text = title

            # 设置图例
            chart.has_legend = has_legend
            if has_legend and legend_position:
                pos_lower = legend_position.lower()
                if pos_lower in PPTOperations.LEGEND_POSITION_MAP:
                    pos_name = PPTOperations.LEGEND_POSITION_MAP[pos_lower]
                    chart.legend.position = getattr(XL_LEGEND_POSITION, pos_name)
                chart.legend.include_in_layout = False

            # 设置数据标签
            if has_data_labels and len(chart.plots) > 0:
                plot = chart.plots[0]
                plot.has_data_labels = True
                if data_label_position:
                    pos_lower = data_label_position.lower()
                    if pos_lower in PPTOperations.LABEL_POSITION_MAP:
                        pos_name = PPTOperations.LABEL_POSITION_MAP[pos_lower]
                        try:
                            plot.data_labels.position = getattr(XL_LABEL_POSITION, pos_name)
                        except (AttributeError, TypeError):
                            pass  # 某些图表类型不支持特定标签位置

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "chart_type": chart_type,
                "series_count": len(series),
                "categories_count": len(categories),
                "message": f"已在第 {slide} 张幻灯片添加 {chart_type} 图表，包含 {len(series)} 个系列"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("ppt_chart_modify",
                   required_params=['path', 'slide', 'chart_index'],
                   optional_params=['title', 'has_legend', 'legend_position',
                                    'has_data_labels', 'data_label_position', 'style'])
    @staticmethod
    def chart_modify(path: str,
                     slide: int,
                     chart_index: int,
                     title: Optional[str] = None,
                     has_legend: Optional[bool] = None,
                     legend_position: Optional[str] = None,
                     has_data_labels: Optional[bool] = None,
                     data_label_position: Optional[str] = None,
                     style: Optional[int] = None) -> Dict[str, Any]:
        """
        修改幻灯片中的图表属性

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            chart_index: 图表索引（从0开始，按添加顺序）
            title: 新的图表标题
            has_legend: 是否显示图例
            legend_position: 图例位置 (right, left, top, bottom, corner)
            has_data_labels: 是否显示数据标签
            data_label_position: 数据标签位置
            style: 图表样式编号（1-48）

        Returns:
            {success, path, slide, chart_index, modifications, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation
            from pptx.enum.chart import XL_LEGEND_POSITION, XL_LABEL_POSITION

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            # 查找图表
            charts = []
            for shape in target_slide.shapes:
                if shape.has_chart:
                    charts.append(shape.chart)

            if not charts:
                return {"success": False, "error": f"第 {slide} 张幻灯片没有图表"}

            if chart_index < 0 or chart_index >= len(charts):
                return {"success": False, "error": f"图表索引超出范围 (0-{len(charts)-1})"}

            chart = charts[chart_index]
            modifications = []

            # 修改标题
            if title is not None:
                chart.has_title = True
                chart.chart_title.text_frame.text = title
                modifications.append("title")

            # 修改图例
            if has_legend is not None:
                chart.has_legend = has_legend
                modifications.append("legend")

            if legend_position is not None and chart.has_legend:
                pos_lower = legend_position.lower()
                if pos_lower in PPTOperations.LEGEND_POSITION_MAP:
                    pos_name = PPTOperations.LEGEND_POSITION_MAP[pos_lower]
                    chart.legend.position = getattr(XL_LEGEND_POSITION, pos_name)
                    modifications.append("legend_position")

            # 修改数据标签
            if has_data_labels is not None and len(chart.plots) > 0:
                chart.plots[0].has_data_labels = has_data_labels
                modifications.append("data_labels")

            if data_label_position is not None and len(chart.plots) > 0:
                pos_lower = data_label_position.lower()
                if pos_lower in PPTOperations.LABEL_POSITION_MAP:
                    pos_name = PPTOperations.LABEL_POSITION_MAP[pos_lower]
                    try:
                        chart.plots[0].data_labels.position = getattr(XL_LABEL_POSITION, pos_name)
                        modifications.append("data_label_position")
                    except (AttributeError, TypeError):
                        pass

            # 修改样式
            if style is not None:
                if 1 <= style <= 48:
                    chart.chart_style = style
                    modifications.append("style")

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "chart_index": chart_index,
                "modifications": modifications,
                "message": f"已修改图表属性: {', '.join(modifications) if modifications else '无修改'}"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("chart_get_data",
                   required_params=['path', 'slide', 'chart_index'],
                   optional_params=[])
    @staticmethod
    def chart_get_data(path: str,
                       slide: int,
                       chart_index: int) -> Dict[str, Any]:
        """
        获取幻灯片中图表的数据

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            chart_index: 图表索引（从0开始）

        Returns:
            {success, path, slide, chart_index, chart_type, title, categories, series, has_legend, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            # 查找图表
            charts = []
            chart_shapes = []
            for shape in target_slide.shapes:
                if shape.has_chart:
                    charts.append(shape.chart)
                    chart_shapes.append(shape)

            if not charts:
                return {"success": False, "error": f"第 {slide} 张幻灯片没有图表"}

            if chart_index < 0 or chart_index >= len(charts):
                return {"success": False, "error": f"图表索引超出范围 (0-{len(charts)-1})"}

            chart = charts[chart_index]
            chart_shape = chart_shapes[chart_index]

            # 获取图表类型
            chart_type_str = str(chart.chart_type).split('.')[-1] if chart.chart_type else "UNKNOWN"

            # 获取标题
            chart_title = ""
            if chart.has_title:
                try:
                    chart_title = chart.chart_title.text_frame.text
                except (AttributeError, ValueError):
                    pass

            # 获取分类（对于分类图表）
            categories = []
            try:
                if chart.plots and len(chart.plots) > 0:
                    plot = chart.plots[0]
                    if hasattr(plot, 'categories') and plot.categories:
                        categories = [str(cat) for cat in plot.categories]
            except (AttributeError, TypeError, ValueError):
                pass

            # 获取系列数据
            series_data = []
            try:
                for series in chart.series:
                    s_info = {
                        "name": series.name if series.name else "Series",
                        "values": []
                    }
                    try:
                        s_info["values"] = list(series.values) if series.values else []
                    except (TypeError, ValueError):
                        pass
                    series_data.append(s_info)
            except (AttributeError, TypeError):
                pass

            # 获取位置信息
            position = {
                "left": chart_shape.left,
                "top": chart_shape.top,
                "width": chart_shape.width,
                "height": chart_shape.height
            }

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "chart_index": chart_index,
                "chart_type": chart_type_str,
                "title": chart_title,
                "categories": categories,
                "series": series_data,
                "series_count": len(series_data),
                "has_legend": chart.has_legend,
                "position": position,
                "message": f"成功获取图表数据: {chart_type_str}，{len(series_data)} 个系列"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("chart_list",
                   required_params=['path', 'slide'],
                   optional_params=[])
    @staticmethod
    def chart_list(path: str,
                   slide: int) -> Dict[str, Any]:
        """
        列出幻灯片中的所有图表

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）

        Returns:
            {success, path, slide, charts, chart_count, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            # 查找所有图表
            charts_info = []
            idx = 0
            for shape in target_slide.shapes:
                if shape.has_chart:
                    chart = shape.chart
                    chart_type_str = str(chart.chart_type).split('.')[-1] if chart.chart_type else "UNKNOWN"

                    chart_title = ""
                    if chart.has_title:
                        try:
                            chart_title = chart.chart_title.text_frame.text
                        except (AttributeError, ValueError):
                            pass

                    series_count = 0
                    try:
                        series_count = len(list(chart.series))
                    except (AttributeError, TypeError):
                        pass

                    charts_info.append({
                        "index": idx,
                        "chart_type": chart_type_str,
                        "title": chart_title,
                        "series_count": series_count,
                        "has_legend": chart.has_legend,
                        "position": {
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height
                        }
                    })
                    idx += 1

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "charts": charts_info,
                "chart_count": len(charts_info),
                "message": f"第 {slide} 张幻灯片共有 {len(charts_info)} 个图表"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    @register_tool("chart_delete",
                   required_params=['path', 'slide', 'chart_index'],
                   optional_params=[])
    @staticmethod
    def chart_delete(path: str,
                     slide: int,
                     chart_index: int) -> Dict[str, Any]:
        """
        删除幻灯片中的图表

        Args:
            path: PPT文件路径
            slide: 幻灯片索引（从1开始）
            chart_index: 图表索引（从0开始）

        Returns:
            {success, path, slide, chart_index, message}
        """
        try:
            if not check_import("pptx"):
                return {"success": False, "error": "需要安装python-pptx: pip install python-pptx"}

            from pptx import Presentation

            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}

            prs = Presentation(path)

            if slide < 1 or slide > len(prs.slides):
                return {"success": False, "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})"}

            target_slide = prs.slides[slide - 1]

            # 查找图表形状
            chart_shapes = []
            for shape in target_slide.shapes:
                if shape.has_chart:
                    chart_shapes.append(shape)

            if not chart_shapes:
                return {"success": False, "error": f"第 {slide} 张幻灯片没有图表"}

            if chart_index < 0 or chart_index >= len(chart_shapes):
                return {"success": False, "error": f"图表索引超出范围 (0-{len(chart_shapes)-1})"}

            # 获取要删除的图表信息
            chart_shape = chart_shapes[chart_index]
            chart_type_str = str(chart_shape.chart.chart_type).split('.')[-1] if chart_shape.chart.chart_type else "UNKNOWN"

            # 删除图表形状
            sp = chart_shape._element
            sp.getparent().remove(sp)

            prs.save(path)

            return {
                "success": True,
                "path": path,
                "slide": slide,
                "chart_index": chart_index,
                "deleted_chart_type": chart_type_str,
                "remaining_charts": len(chart_shapes) - 1,
                "message": f"已删除第 {slide} 张幻灯片的第 {chart_index} 个图表 ({chart_type_str})"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}
