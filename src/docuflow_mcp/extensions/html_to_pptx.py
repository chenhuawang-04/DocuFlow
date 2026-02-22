# -*- coding: utf-8 -*-
"""
HTML to PPTX Converter
将HTML网页(div+p标签,绝对定位,内联样式)转换为PowerPoint演示文稿

支持的CSS属性:
- 位置: position:absolute, left, top, right, bottom
- 大小: width, height
- 背景: background-color, background (linear-gradient, radial-gradient), rgba透明度
- 边框: border-radius (圆角)
- 文本: font-size, font-weight, font-family, color (含rgba透明度), text-align
"""

import os
import re
from typing import Dict, Any, Optional, Tuple

from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from docuflow_mcp.core.registry import register_tool


class CSSParser:
    """CSS样式解析器"""

    @staticmethod
    def parse_inline_style(style_str: str) -> Dict[str, str]:
        """解析内联style属性"""
        if not style_str:
            return {}
        styles = {}
        for decl in style_str.split(';'):
            decl = decl.strip()
            if ':' in decl:
                prop, value = decl.split(':', 1)
                styles[prop.strip().lower()] = value.strip()
        return styles

    @staticmethod
    def parse_length(value: str, reference: float = 1920) -> float:
        """解析CSS长度值,返回像素值"""
        if not value:
            return 0
        value = value.strip().lower()
        match = re.match(r'^(-?\d*\.?\d+)(px|pt|in|cm|%)?$', value)
        if match:
            num = float(match.group(1))
            unit = match.group(2) or 'px'
            if unit == 'px':
                return num
            elif unit == 'pt':
                return num * 1.333
            elif unit == 'in':
                return num * 96
            elif unit == 'cm':
                return num * 37.795
            elif unit == '%':
                return num / 100 * reference
        return 0

    @staticmethod
    def parse_color(value: str) -> Optional[Tuple[int, int, int, float]]:
        """解析颜色,返回(r, g, b, alpha)"""
        if not value:
            return None
        value = value.strip().lower()

        # #hex
        if value.startswith('#'):
            h = value[1:]
            if len(h) == 3:
                h = ''.join([c * 2 for c in h])
            if len(h) == 6:
                return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16), 1.0)

        # rgba(r, g, b, a)
        m = re.match(r'rgba\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*,\s*([\d.]+)\s*\)', value)
        if m:
            return (int(m.group(1)), int(m.group(2)), int(m.group(3)), float(m.group(4)))

        # rgb(r, g, b)
        m = re.match(r'rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', value)
        if m:
            return (int(m.group(1)), int(m.group(2)), int(m.group(3)), 1.0)

        # 颜色名
        colors = {
            'white': (255, 255, 255, 1.0),
            'black': (0, 0, 0, 1.0),
            'red': (255, 0, 0, 1.0),
            'green': (0, 128, 0, 1.0),
            'blue': (0, 0, 255, 1.0),
            'yellow': (255, 255, 0, 1.0),
            'cyan': (0, 255, 255, 1.0),
            'magenta': (255, 0, 255, 1.0),
            'gray': (128, 128, 128, 1.0),
            'grey': (128, 128, 128, 1.0),
            'orange': (255, 165, 0, 1.0),
            'purple': (128, 0, 128, 1.0),
            'pink': (255, 192, 203, 1.0),
            'transparent': None,
        }
        if value in colors:
            return colors[value]

        return None

    @staticmethod
    def parse_gradient(value: str) -> Optional[Dict[str, Any]]:
        """解析CSS渐变,支持linear-gradient和radial-gradient"""
        if not value:
            return None
        m = re.match(r'(linear|radial)-gradient\s*\(\s*(.+)\s*\)', value, re.IGNORECASE)
        if not m:
            return None

        grad_type = m.group(1).lower()
        content = m.group(2)
        parts = [p.strip() for p in re.split(r',(?![^()]*\))', content)]
        angle = 180
        stops = []

        for i, part in enumerate(parts):
            if i == 0:
                # Try to parse angle (linear only)
                am = re.match(r'(\d+)deg', part)
                if am:
                    angle = int(am.group(1))
                    continue
                # Skip non-color directives (circle, ellipse, at center, to right, etc.)
                if grad_type == 'radial' and not re.match(r'(rgba?\s*\(|#)', part):
                    continue

            # Extract color: match rgba(...) or rgb(...) or #hex or named color
            cm = re.match(r'(rgba?\s*\([^)]+\)|#[0-9a-fA-F]{3,8}|\w+)', part)
            if not cm:
                continue
            color_str = cm.group(1)
            color = CSSParser.parse_color(color_str)
            if not color:
                continue

            r, g, b, alpha = color

            # Extract position percentage
            pos_match = re.search(r'([\d.]+)%', part[cm.end():])
            position = float(pos_match.group(1)) if pos_match else None

            stops.append((r, g, b, alpha, position))

        if stops:
            return {
                'type': grad_type,
                'angle': angle,
                'stops': stops,
            }
        return None


class HTMLToPPTXConverter:
    """HTML到PPTX转换器"""

    # 默认幻灯片尺寸
    SLIDE_W = 1920
    SLIDE_H = 1080
    PPT_W = 13.333
    PPT_H = 7.5

    def __init__(self, slide_width: int = 1920, slide_height: int = 1080):
        """
        初始化转换器

        Args:
            slide_width: HTML画布宽度(像素)
            slide_height: HTML画布高度(像素)
        """
        self.slide_w = slide_width
        self.slide_h = slide_height
        self.ppt_w = self.PPT_W
        self.ppt_h = self.PPT_H
        self.scale_x = self.ppt_w / self.slide_w
        self.scale_y = self.ppt_h / self.slide_h
        self.font_scale = self.ppt_w * 72 / self.slide_w  # px -> pt
        self.css = CSSParser()

    def px_to_in(self, px: float, axis: str = 'x') -> float:
        """像素转英寸"""
        scale = self.scale_x if axis == 'x' else self.scale_y
        return px * scale

    def set_rounded_corners(self, shape, radius_inches: float):
        """设置形状圆角"""
        spPr = shape._element.spPr
        prstGeom = spPr.prstGeom
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            min_dim = min(shape.width.inches, shape.height.inches)
            if min_dim > 0:
                ratio = min(radius_inches / min_dim * 50000, 50000)
                for child in list(avLst):
                    avLst.remove(child)
                gd = etree.SubElement(avLst, qn('a:gd'))
                gd.set('name', 'adj')
                gd.set('fmla', f'val {int(ratio)}')

    def apply_fill(self, shape, styles: Dict[str, str]):
        """应用形状填充"""
        bg_value = styles.get('background', '') or styles.get('background-color', '')
        grad = self.css.parse_gradient(bg_value)
        if grad and grad['stops']:
            self._apply_gradient_xml(shape, grad)
        else:
            color = self.css.parse_color(bg_value)
            if color:
                r, g, b, alpha = color
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(r, g, b)
                if alpha < 1.0:
                    spPr = shape._element.spPr
                    solidFill = spPr.find(qn('a:solidFill'))
                    if solidFill is not None:
                        srgbClr = solidFill.find(qn('a:srgbClr'))
                        if srgbClr is not None:
                            alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
                            alpha_elem.set('val', str(int(alpha * 100000)))
            else:
                # No parseable color/gradient → transparent fill
                shape.fill.background()

    # PPT渲染引擎对径向渐变的alpha表现远强于浏览器,
    # 用幂次曲线压缩低alpha值以补偿视觉差异
    RADIAL_ALPHA_POWER = 2.0

    def _apply_gradient_xml(self, shape, grad: Dict[str, Any]):
        """通过OOXML直接创建渐变填充,支持alpha透明度和多停止点"""
        spPr = shape._element.spPr

        # Remove existing fill elements
        for tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:pattFill', 'a:blipFill'):
            for old in spPr.findall(qn(tag)):
                spPr.remove(old)

        # Build gradFill as detached element first
        gradFill = etree.Element(qn('a:gradFill'))
        gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

        stops = grad['stops']
        is_radial = grad.get('type') == 'radial'

        for i, stop in enumerate(stops):
            r, g, b, alpha, position = stop

            # 径向渐变alpha补偿: alpha^2 把 0.08 → 0.0064, 0.5 → 0.25, 1.0 → 1.0
            if is_radial and alpha < 1.0:
                alpha = alpha ** self.RADIAL_ALPHA_POWER
            # Auto-assign positions if missing
            if position is None:
                if len(stops) == 1:
                    position = 0
                else:
                    position = i / (len(stops) - 1) * 100
            pos_val = int(position * 1000)  # percentage * 1000 → 0-100000

            gs = etree.SubElement(gsLst, qn('a:gs'))
            gs.set('pos', str(pos_val))

            srgbClr = etree.SubElement(gs, qn('a:srgbClr'))
            srgbClr.set('val', f'{r:02X}{g:02X}{b:02X}')

            if alpha < 1.0:
                alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha_elem.set('val', str(int(alpha * 100000)))

        if is_radial:
            path_elem = etree.SubElement(gradFill, qn('a:path'))
            path_elem.set('path', 'circle')
            fillToRect = etree.SubElement(path_elem, qn('a:fillToRect'))
            fillToRect.set('l', '50000')
            fillToRect.set('t', '50000')
            fillToRect.set('r', '50000')
            fillToRect.set('b', '50000')
        else:
            css_angle = grad.get('angle', 180)
            ooxml_angle = ((css_angle - 90 + 360) % 360) * 60000
            lin = etree.SubElement(gradFill, qn('a:lin'))
            lin.set('ang', str(ooxml_angle))
            lin.set('scaled', '1')

        # Insert in correct OOXML order: fill must come before a:ln
        ln = spPr.find(qn('a:ln'))
        if ln is not None:
            ln.addprevious(gradFill)
        else:
            spPr.append(gradFill)

    def apply_text_style(self, run, paragraph, styles: Dict[str, str]):
        """应用文本样式"""
        font = run.font

        # 字体大小
        fs = self.css.parse_length(styles.get('font-size', '16'))
        if fs > 0:
            font.size = Pt(fs * self.font_scale)

        # 颜色
        color_str = styles.get('color', '')
        color = self.css.parse_color(color_str)
        if color:
            r, g, b, alpha = color
            font.color.rgb = RGBColor(r, g, b)
            if alpha < 1.0:
                rPr = run._r.get_or_add_rPr()
                solidFill = rPr.find(qn('a:solidFill'))
                if solidFill is not None:
                    srgbClr = solidFill.find(qn('a:srgbClr'))
                    if srgbClr is not None:
                        for old_alpha in srgbClr.findall(qn('a:alpha')):
                            srgbClr.remove(old_alpha)
                        alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
                        alpha_elem.set('val', str(int(alpha * 100000)))

        # 粗体
        fw = styles.get('font-weight', '')
        if fw in ('bold', '700', '800', '900'):
            font.bold = True

        # 字体
        ff = styles.get('font-family', '')
        if ff:
            font.name = ff.split(',')[0].strip().strip("'\"")

        # 对齐
        ta = styles.get('text-align', '')
        if ta == 'center':
            paragraph.alignment = PP_ALIGN.CENTER
        elif ta == 'right':
            paragraph.alignment = PP_ALIGN.RIGHT

    def _convert_slide(self, prs: Presentation, html_content: str) -> int:
        """
        将单个HTML转换为Presentation中的一张幻灯片

        Args:
            prs: 已有的Presentation对象
            html_content: HTML内容字符串

        Returns:
            添加的元素数量, 失败返回-1
        """
        soup = BeautifulSoup(html_content, 'html.parser')
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 找主容器
        body = soup.find('body')
        main = None
        for child in body.children:
            if isinstance(child, Tag):
                s = self.css.parse_inline_style(child.get('style', ''))
                if s.get('width') and s.get('height'):
                    # 检测实际尺寸
                    w = self.css.parse_length(s['width'])
                    h = self.css.parse_length(s['height'])
                    if w > 0 and h > 0:
                        self.slide_w = w
                        self.slide_h = h
                        self.scale_x = self.ppt_w / self.slide_w
                        self.scale_y = self.ppt_h / self.slide_h
                        self.font_scale = self.ppt_w * 72 / self.slide_w
                    main = child
                    break

        if not main:
            return -1

        # 背景
        ms = self.css.parse_inline_style(main.get('style', ''))
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(self.ppt_w), Inches(self.ppt_h)
        )
        bg.line.fill.background()
        self.apply_fill(bg, ms)

        # 处理子元素
        count = 0
        for elem in main.descendants:
            if not isinstance(elem, Tag):
                continue
            s = self.css.parse_inline_style(elem.get('style', ''))
            if s.get('position') != 'absolute':
                continue

            # 解析位置和大小
            left = self.css.parse_length(s.get('left', '0'))
            top = self.css.parse_length(s.get('top', '0'))
            w = self.css.parse_length(s.get('width', '100'))
            h = self.css.parse_length(s.get('height', '50'))

            # 处理 bottom/right
            if 'bottom' in s and 'top' not in s:
                bottom = self.css.parse_length(s['bottom'])
                top = self.slide_h - bottom - h
            if 'right' in s and 'left' not in s:
                right = self.css.parse_length(s['right'])
                left = self.slide_w - right - w

            # 转换为英寸
            L = Inches(self.px_to_in(left, 'x'))
            T = Inches(self.px_to_in(top, 'y'))
            W = Inches(self.px_to_in(w, 'x'))
            H = Inches(self.px_to_in(h, 'y'))

            # 创建形状
            if elem.name == 'div':
                radius_px = self.css.parse_length(s.get('border-radius', '0'))
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, L, T, W, H)
                if radius_px > 0:
                    self.set_rounded_corners(shape, self.px_to_in(radius_px, 'x'))
                shape.line.fill.background()
                self.apply_fill(shape, s)
                count += 1

            elif elem.name == 'p':
                text = elem.get_text(strip=True)
                if text:
                    has_explicit_width = 'width' in s

                    if not has_explicit_width:
                        # Auto-estimate width: CJK ~1em, Latin ~0.6em
                        fs_px = self.css.parse_length(s.get('font-size', '16'))
                        char_w = sum(1.0 if ord(c) > 0x2000 else 0.6 for c in text)
                        W = Inches(self.px_to_in(char_w * fs_px * 1.3, 'x'))

                    tb = slide.shapes.add_textbox(L, T, W, H)
                    tf = tb.text_frame
                    tf.word_wrap = has_explicit_width
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = text
                    self.apply_text_style(run, p, s)
                    count += 1

        return count

    def convert(self, html_content: str, output_path: str) -> Dict[str, Any]:
        """
        将单个HTML转换为PPTX文件

        Args:
            html_content: HTML内容字符串
            output_path: 输出PPTX文件路径

        Returns:
            转换结果
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(self.ppt_w)
            prs.slide_height = Inches(self.ppt_h)

            count = self._convert_slide(prs, html_content)
            if count < 0:
                return {'success': False, 'error': 'No main container found'}

            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)

            prs.save(output_path)

            return {
                'success': True,
                'output_path': output_path,
                'slide_size': f'{int(self.slide_w)}x{int(self.slide_h)}',
                'elements_count': count,
                'message': f'Successfully converted HTML to PPTX with {count} elements'
            }

        except Exception as e:
            import traceback
            return {
                'success': False,
                'error': str(e),
                'traceback': traceback.format_exc()
            }

    def convert_multi(self, html_contents: list, output_path: str) -> Dict[str, Any]:
        """
        将多个HTML转换为一个多页PPTX文件

        Args:
            html_contents: HTML内容字符串列表
            output_path: 输出PPTX文件路径

        Returns:
            转换结果
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(self.ppt_w)
            prs.slide_height = Inches(self.ppt_h)

            total = 0
            for html in html_contents:
                count = self._convert_slide(prs, html)
                if count >= 0:
                    total += count

            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)

            prs.save(output_path)

            return {
                'success': True,
                'output_path': output_path,
                'slide_size': f'{int(self.slide_w)}x{int(self.slide_h)}',
                'total_slides': len(prs.slides),
                'elements_count': total,
                'message': f'Successfully converted {len(prs.slides)} slides with {total} elements'
            }

        except Exception as e:
            import traceback
            return {
                'success': False,
                'error': str(e),
                'traceback': traceback.format_exc()
            }


class HTMLToPPTXOperations:
    """HTML转PPTX MCP工具"""

    @register_tool("html_to_pptx_convert",
                   required_params=['html_source', 'output_path'],
                   optional_params=['base_path'])
    @staticmethod
    def convert(html_source: str, output_path: str, base_path: str = None) -> Dict[str, Any]:
        """
        将HTML转换为PPTX

        Args:
            html_source: HTML文件路径或HTML内容字符串
            output_path: 输出PPTX文件路径
            base_path: 基础路径(用于解析相对图片路径,可选)

        Returns:
            转换结果
        """
        # 判断是文件路径还是HTML内容
        if os.path.isfile(html_source):
            with open(html_source, 'r', encoding='utf-8') as f:
                html_content = f.read()
        else:
            html_content = html_source

        converter = HTMLToPPTXConverter()
        return converter.convert(html_content, output_path)

    @register_tool("html_to_pptx_convert_multi",
                   required_params=['html_sources', 'output_path'],
                   optional_params=[])
    @staticmethod
    def convert_multi(html_sources: list, output_path: str) -> Dict[str, Any]:
        """
        将多个HTML转换为一个多页PPTX文件

        Args:
            html_sources: HTML文件路径或HTML内容字符串的列表
            output_path: 输出PPTX文件路径

        Returns:
            转换结果
        """
        html_contents = []
        for src in html_sources:
            if os.path.isfile(src):
                with open(src, 'r', encoding='utf-8') as f:
                    html_contents.append(f.read())
            else:
                html_contents.append(src)

        converter = HTMLToPPTXConverter()
        return converter.convert_multi(html_contents, output_path)

    @register_tool("html_to_pptx_status", required_params=[], optional_params=[])
    @staticmethod
    def get_status() -> Dict[str, Any]:
        """获取HTML转PPTX模块状态"""
        try:
            from bs4 import BeautifulSoup
            bs4_ok = True
        except ImportError:
            bs4_ok = False

        try:
            from pptx import Presentation
            pptx_ok = True
        except ImportError:
            pptx_ok = False

        return {
            'available': bs4_ok and pptx_ok,
            'dependencies': {
                'beautifulsoup4': bs4_ok,
                'python-pptx': pptx_ok
            },
            'default_slide_size': '1920x1080 (16:9)',
            'supported_css': [
                'position: absolute',
                'left, top, right, bottom',
                'width, height',
                'background-color (with rgba)',
                'background: linear-gradient()',
                'background: radial-gradient()',
                'border-radius',
                'font-size, font-weight, font-family',
                'color (with rgba transparency)',
                'text-align'
            ],
            'supported_elements': ['div', 'p'],
            'notes': 'HTML must use absolute positioning for accurate conversion'
        }
