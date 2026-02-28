"""
DocuFlow ImageGen - AI图片生成模块

通过AI API生成图片，支持在PPT中插入AI生成的图片
基于 aiphotogen 项目改造
"""
import base64
import datetime as dt
import json
import os
import re
import ssl
import uuid
from pathlib import Path
from typing import Any, Dict, Optional, Tuple
from urllib.error import HTTPError, URLError
from urllib.parse import urljoin
from urllib.request import (
    HTTPHandler,
    HTTPRedirectHandler,
    HTTPSHandler,
    Request,
    build_opener,
    urlopen,
)

from ..core.registry import register_tool


# 默认配置
DEFAULT_API_URL = "https://ai.com/v1/chat/completions"
DEFAULT_MODEL = "gpt-4o-mini"
DEFAULT_TIMEOUT = 120
DEFAULT_OUTPUT_DIR = "generated_images"

# 配置文件路径
CONFIG_FILE = Path(__file__).resolve().parent.parent.parent.parent / "image_gen_config.json"


class ImageExtractionError(RuntimeError):
    """图片提取错误"""
    pass


def _load_config() -> Dict[str, Any]:
    """加载配置文件"""
    if CONFIG_FILE.is_file():
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except (json.JSONDecodeError, IOError):
            pass
    return {}


def _get_api_key(config: Dict[str, Any]) -> str:
    """获取API密钥（仅从环境变量读取，不从配置文件读取以防泄露）"""
    for var in ("AI_API_KEY", "OPENAI_API_KEY"):
        value = os.getenv(var)
        if value:
            return value
    raise RuntimeError(
        "缺少API密钥。请设置环境变量 AI_API_KEY 或 OPENAI_API_KEY。"
    )


def _build_payload(prompt: str, model: str, extra: Dict[str, Any] = None) -> Dict[str, Any]:
    """构建请求payload"""
    payload = {
        "model": model,
        "messages": [{"role": "user", "content": prompt}],
    }
    if extra:
        payload.update(extra)
    return payload


def _request_chat_completion(
    api_url: str,
    api_key: str,
    payload: Dict[str, Any],
    timeout: int,
    ssl_context: Optional[ssl.SSLContext],
    redirects_remaining: int = 3,
) -> Tuple[Dict[str, Any], str]:
    """发送API请求"""
    body = json.dumps(payload).encode("utf-8")
    request = Request(
        api_url,
        data=body,
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        },
        method="POST",
    )

    class NoRedirectHandler(HTTPRedirectHandler):
        def redirect_request(self, req, fp, code, msg, headers, newurl):
            return None

    opener = build_opener(
        NoRedirectHandler(),
        HTTPHandler(),
        HTTPSHandler(context=ssl_context) if ssl_context else HTTPSHandler(),
    )

    try:
        with opener.open(request, timeout=timeout) as response:
            raw = response.read()
    except HTTPError as exc:
        if exc.code in (301, 302, 303, 307, 308) and redirects_remaining > 0:
            location = exc.headers.get("Location")
            if location:
                redirected_url = urljoin(api_url, location)
                # 安全检查：防止重定向到不同域名泄漏 API Key
                from urllib.parse import urlparse
                orig_host = urlparse(api_url).hostname
                redir_host = urlparse(redirected_url).hostname
                if orig_host != redir_host:
                    raise RuntimeError(
                        f"API 重定向到不同域名 ({redir_host})，已阻止以防 API Key 泄漏"
                    ) from exc
                return _request_chat_completion(
                    redirected_url, api_key, payload, timeout, ssl_context, redirects_remaining - 1
                )
        detail = exc.read().decode("utf-8", errors="replace")
        raise RuntimeError(f"API请求失败 ({exc.code}): {detail}") from exc
    except URLError as exc:
        raise RuntimeError(f"API请求失败: {exc.reason}") from exc

    raw_text = raw.decode("utf-8", errors="replace")
    try:
        data = json.loads(raw_text)
    except json.JSONDecodeError as exc:
        raise RuntimeError("API响应不是有效的JSON") from exc
    if not isinstance(data, dict):
        raise RuntimeError("API响应JSON不是对象")
    return data, raw_text


def _extract_from_content(content: Any) -> Optional[Tuple[str, str]]:
    """从内容中提取图片数据"""
    if not content:
        return None
    if isinstance(content, str):
        stripped = content.strip()
        if stripped.startswith("http://") or stripped.startswith("https://"):
            return "url", stripped
        if stripped.startswith("data:image/"):
            return "data_url", stripped
        match = re.search(r"!\[[^\]]*]\((https?://[^)]+)\)", stripped)
        if match:
            return "url", match.group(1)
        url_match = re.search(r"https?://\S+", stripped)
        if url_match:
            return "url", url_match.group(0).rstrip(").,")
        if stripped.startswith("{"):
            try:
                parsed = json.loads(stripped)
            except json.JSONDecodeError:
                return None
            return _extract_from_content(parsed)
        return None
    if isinstance(content, dict):
        if "b64_json" in content:
            return "base64", content["b64_json"]
        if "image_base64" in content:
            return "base64", content["image_base64"]
        if "url" in content:
            return "url", content["url"]
        if content.get("type") in ("image_url", "output_image"):
            image_url = content.get("image_url", {})
            if isinstance(image_url, dict) and "url" in image_url:
                return "url", image_url["url"]
            if isinstance(image_url, str):
                return "url", image_url
        if "image_url" in content and isinstance(content["image_url"], dict):
            url_value = content["image_url"].get("url")
            if isinstance(url_value, str):
                return "url", url_value
        return None
    if isinstance(content, list):
        for item in content:
            found = _extract_from_content(item)
            if found:
                return found
    return None


def _extract_image_payload(response: Dict[str, Any]) -> Tuple[str, str]:
    """从API响应中提取图片数据"""
    # 检查 data 字段
    for item in response.get("data", []) or []:
        if isinstance(item, dict):
            if "b64_json" in item:
                return "base64", item["b64_json"]
            if "image_base64" in item:
                return "base64", item["image_base64"]
            if "url" in item:
                return "url", item["url"]

    # 检查 choices 字段
    choices = response.get("choices", []) or []
    if isinstance(choices, list):
        for choice in choices:
            if not isinstance(choice, dict):
                continue
            message = choice.get("message", {})
            found = _extract_from_content(message.get("content"))
            if found:
                return found

    # 检查 output 字段
    output = response.get("output")
    found = _extract_from_content(output)
    if found:
        return found

    raise ImageExtractionError("响应中未找到图片数据")


def _decode_image(
    source_type: str,
    payload: str,
    timeout: int,
    ssl_context: Optional[ssl.SSLContext],
) -> bytes:
    """解码图片数据"""
    if source_type == "base64":
        return base64.b64decode(payload)
    if source_type == "data_url":
        header, _, data = payload.partition(",")
        if not header.startswith("data:image/"):
            raise ImageExtractionError("不支持的data URL格式")
        return base64.b64decode(data)
    if source_type == "url":
        request = Request(payload, headers={"User-Agent": "docuflow-imagegen"})
        try:
            with urlopen(request, timeout=timeout, context=ssl_context) as response:
                return response.read()
        except HTTPError as exc:
            raise ImageExtractionError(f"图片下载失败 ({exc.code})") from exc
        except URLError as exc:
            raise ImageExtractionError(f"图片下载失败: {exc.reason}") from exc
    raise ImageExtractionError(f"不支持的来源类型: {source_type}")


def _detect_image_format(data: bytes) -> Optional[str]:
    """检测图片格式"""
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return "png"
    if data.startswith(b"\xff\xd8"):
        return "jpg"
    if data.startswith(b"GIF87a") or data.startswith(b"GIF89a"):
        return "gif"
    if len(data) > 12 and data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return "webp"
    return None


def _read_png_size(data: bytes) -> Optional[Tuple[int, int]]:
    if len(data) < 24:
        return None
    return int.from_bytes(data[16:20], "big"), int.from_bytes(data[20:24], "big")


def _read_jpeg_size(data: bytes) -> Optional[Tuple[int, int]]:
    if len(data) < 4 or not data.startswith(b"\xff\xd8"):
        return None
    index = 2
    while index < len(data) - 1:
        if data[index] != 0xFF:
            index += 1
            continue
        marker = data[index + 1]
        if marker in (0xD8, 0xD9):
            index += 2
            continue
        if index + 4 >= len(data):
            return None
        length = int.from_bytes(data[index + 2:index + 4], "big")
        if length < 2:
            return None
        if marker in (0xC0, 0xC1, 0xC2, 0xC3, 0xC5, 0xC6, 0xC7, 0xC9, 0xCA, 0xCB, 0xCD, 0xCE, 0xCF):
            if index + 7 >= len(data):
                return None
            height = int.from_bytes(data[index + 5:index + 7], "big")
            width = int.from_bytes(data[index + 7:index + 9], "big")
            return width, height
        index += 2 + length
    return None


def _get_image_dimensions(data: bytes, image_format: Optional[str]) -> Tuple[Optional[int], Optional[int]]:
    """获取图片尺寸"""
    if image_format == "png":
        return _read_png_size(data) or (None, None)
    if image_format == "jpg":
        return _read_jpeg_size(data) or (None, None)
    return None, None


def _save_image(data: bytes, out_dir: str, filename: Optional[str] = None) -> Path:
    """保存图片到文件"""
    target_dir = Path(out_dir)
    target_dir.mkdir(parents=True, exist_ok=True)
    image_format = _detect_image_format(data) or "png"
    if not filename:
        filename = f"ai_image_{dt.datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:8]}.{image_format}"
    elif not filename.endswith(f".{image_format}"):
        filename = f"{filename}.{image_format}"
    path = target_dir / filename
    path.write_bytes(data)
    return path


class ImageGenOperations:
    """AI图片生成操作"""

    @register_tool("image_gen_status",
                   required_params=[],
                   optional_params=[])
    @staticmethod
    def get_status() -> Dict[str, Any]:
        """
        获取图片生成模块状态

        Returns:
            {success, config_exists, api_url, model, features}
        """
        config = _load_config()
        has_api_key = bool(config.get("api_key") or os.getenv("AI_API_KEY") or os.getenv("OPENAI_API_KEY"))

        return {
            "success": True,
            "config_file": str(CONFIG_FILE),
            "config_exists": CONFIG_FILE.is_file(),
            "has_api_key": has_api_key,
            "api_url": config.get("api_url", DEFAULT_API_URL),
            "model": config.get("model", DEFAULT_MODEL),
            "default_output_dir": DEFAULT_OUTPUT_DIR,
            "features": [
                "AI图片生成",
                "自动格式检测",
                "尺寸获取",
                "PPT集成"
            ]
        }

    @register_tool("image_generate",
                   required_params=['prompt'],
                   optional_params=['output_dir', 'filename', 'timeout', 'model', 'api_url'])
    @staticmethod
    def generate(
        prompt: str,
        output_dir: Optional[str] = None,
        filename: Optional[str] = None,
        timeout: Optional[int] = None,
        model: Optional[str] = None,
        api_url: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        使用AI生成图片

        Args:
            prompt: 图片描述提示词，描述你想要生成的图片内容
            output_dir: 输出目录(可选，默认 generated_images)
            filename: 文件名(可选，自动生成)
            timeout: 超时时间(秒)，默认120
            model: 模型名称(可选，使用配置文件中的默认值)
            api_url: API地址(可选，使用配置文件中的默认值)

        Returns:
            {success, path, format, width, height, size_bytes}
        """
        try:
            config = _load_config()

            # 获取配置
            api_key = _get_api_key(config)
            final_api_url = api_url or config.get("api_url") or DEFAULT_API_URL
            final_model = model or config.get("model") or DEFAULT_MODEL
            final_timeout = timeout or DEFAULT_TIMEOUT
            final_output_dir = output_dir or DEFAULT_OUTPUT_DIR

            # 构建请求
            payload = _build_payload(prompt, final_model)
            ssl_context = ssl.create_default_context()

            # 发送请求
            response, raw_text = _request_chat_completion(
                final_api_url, api_key, payload, final_timeout, ssl_context
            )

            # 提取图片
            source_type, payload_data = _extract_image_payload(response)
            image_bytes = _decode_image(source_type, payload_data, final_timeout, ssl_context)

            # 保存图片
            path = _save_image(image_bytes, final_output_dir, filename)

            # 获取图片信息
            image_format = _detect_image_format(image_bytes)
            width, height = _get_image_dimensions(image_bytes, image_format)

            return {
                "success": True,
                "path": str(path.resolve()),
                "format": image_format,
                "width": width,
                "height": height,
                "size_bytes": len(image_bytes),
                "source": source_type,
                "prompt": prompt,
                "model": final_model,
                "message": f"已生成图片: {path}"
            }

        except (RuntimeError, ImageExtractionError, ValueError) as exc:
            return {
                "success": False,
                "error": str(exc)
            }

    @register_tool("image_generate_for_ppt",
                   required_params=['ppt_path', 'slide', 'prompt'],
                   optional_params=['left', 'top', 'width', 'height', 'output_dir', 'timeout', 'model'])
    @staticmethod
    def generate_for_ppt(
        ppt_path: str,
        slide: int,
        prompt: str,
        left: Optional[str] = "1in",
        top: Optional[str] = "1in",
        width: Optional[str] = None,
        height: Optional[str] = None,
        output_dir: Optional[str] = None,
        timeout: Optional[int] = None,
        model: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        生成AI图片并插入到PPT幻灯片中

        Args:
            ppt_path: PPT文件路径
            slide: 幻灯片索引(从1开始)
            prompt: 图片描述提示词
            left: 图片左边距，如 '1in', '2.54cm'
            top: 图片上边距
            width: 图片宽度(可选，保持比例)
            height: 图片高度(可选)
            output_dir: 图片保存目录(可选)
            timeout: 超时时间(秒)
            model: 模型名称(可选)

        Returns:
            {success, image_path, ppt_path, slide, message}
        """
        try:
            # 先生成图片
            gen_result = ImageGenOperations.generate(
                prompt=prompt,
                output_dir=output_dir or DEFAULT_OUTPUT_DIR,
                timeout=timeout,
                model=model
            )

            if not gen_result.get("success"):
                return gen_result

            image_path = gen_result["path"]

            # 检查PPT库
            try:
                from pptx import Presentation
                from pptx.util import Inches, Cm
            except ImportError:
                return {
                    "success": False,
                    "error": "需要安装python-pptx: pip install python-pptx",
                    "image_path": image_path
                }

            # 解析长度
            def parse_length(value):
                if not value:
                    return None
                value = str(value).strip().lower()
                if value.endswith('in'):
                    return Inches(float(value[:-2]))
                elif value.endswith('cm'):
                    return Cm(float(value[:-2]))
                else:
                    try:
                        return Inches(float(value))
                    except ValueError:
                        return Inches(1)

            # 打开PPT
            prs = Presentation(ppt_path)

            if slide < 1 or slide > len(prs.slides):
                return {
                    "success": False,
                    "error": f"幻灯片索引超出范围 (1-{len(prs.slides)})",
                    "image_path": image_path
                }

            target_slide = prs.slides[slide - 1]

            # 添加图片
            left_val = parse_length(left) or Inches(1)
            top_val = parse_length(top) or Inches(1)
            width_val = parse_length(width) if width else None
            height_val = parse_length(height) if height else None

            target_slide.shapes.add_picture(
                image_path,
                left_val,
                top_val,
                width=width_val,
                height=height_val
            )

            # 保存PPT
            prs.save(ppt_path)

            return {
                "success": True,
                "image_path": image_path,
                "image_format": gen_result.get("format"),
                "image_width": gen_result.get("width"),
                "image_height": gen_result.get("height"),
                "ppt_path": ppt_path,
                "slide": slide,
                "position": {"left": left, "top": top},
                "prompt": prompt,
                "message": f"已生成图片并插入到第 {slide} 张幻灯片"
            }

        except Exception as exc:
            return {
                "success": False,
                "error": str(exc)
            }
